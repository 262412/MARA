from __future__ import annotations

import csv
import html
import json
from pathlib import Path
from typing import Any, Callable, cast

from theflow.settings import settings as flowsettings
from theflow.utils.modules import import_dotted_string

MEDIA_EXPORT_FORMATS = {"mp3", "mp4"}
MEDIA_EXPORT_ADAPTER_SETTING = "KH_MARA_ARTIFACT_MEDIA_EXPORT_ADAPTER"
MediaExportAdapter = Callable[[dict[str, Any], str, Path], str | Path]


def export_artifact_to_path(
    artifact: dict[str, Any],
    *,
    export_format: str,
    output_path: str | Path,
    media_export_adapter: MediaExportAdapter | None = None,
) -> Path:
    target = Path(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)
    normalized_format = str(export_format or "").lower().strip()
    if normalized_format == "json":
        target.write_text(
            json.dumps(artifact, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
    elif normalized_format == "html":
        target.write_text(_render_html(artifact), encoding="utf-8")
    elif normalized_format == "csv":
        _write_csv(artifact, target)
    elif normalized_format in {"md", "markdown"}:
        target.write_text(render_artifact_markdown(artifact), encoding="utf-8")
    elif normalized_format == "svg":
        target.write_text(_render_svg(artifact), encoding="utf-8")
    elif normalized_format == "pptx":
        _write_pptx(artifact, target)
    elif normalized_format in MEDIA_EXPORT_FORMATS:
        adapter = media_export_adapter or configured_media_export_adapter()
        if adapter is not None:
            return Path(adapter(artifact, normalized_format, target))
        raise ValueError(
            f"Artifact export format '{export_format}' requires a configured "
            "media export adapter."
        )
    else:
        raise ValueError(f"Unsupported artifact export format '{export_format}'.")
    return target


def configured_media_export_adapter() -> MediaExportAdapter | None:
    dotted_path = str(
        getattr(flowsettings, MEDIA_EXPORT_ADAPTER_SETTING, "") or ""
    ).strip()
    if not dotted_path:
        return None
    adapter = import_dotted_string(dotted_path, safe=False)
    if not callable(adapter):
        raise ValueError(f"{MEDIA_EXPORT_ADAPTER_SETTING} must resolve to a callable.")
    return cast(MediaExportAdapter, adapter)


def render_artifact_markdown(artifact: dict[str, Any]) -> str:
    title = str(artifact.get("title") or artifact.get("type") or "Artifact")
    payload = artifact.get("payload")
    lines = [f"# {title}", "", f"Type: {artifact.get('type', '')}", ""]
    if isinstance(payload, dict):
        lines.extend(_payload_markdown_lines(payload))
    else:
        lines.append(str(payload or ""))
    citations = artifact.get("citations") or []
    if citations:
        lines.extend(["", "## Citations", ""])
        lines.extend(f"- {json.dumps(item, ensure_ascii=False)}" for item in citations)
    return "\n".join(lines).rstrip() + "\n"


def _payload_markdown_lines(payload: dict[str, Any]) -> list[str]:
    lines: list[str] = []
    for key, value in payload.items():
        heading = str(key).replace("_", " ").title()
        lines.extend([f"## {heading}", ""])
        if isinstance(value, list):
            lines.extend(f"- {item}" for item in value)
        elif isinstance(value, dict):
            lines.append(json.dumps(value, ensure_ascii=False, indent=2))
        else:
            lines.append(str(value or ""))
        lines.append("")
    return lines


def _render_html(artifact: dict[str, Any]) -> str:
    markdown = render_artifact_markdown(artifact)
    body = "<br>\n".join(html.escape(line) for line in markdown.splitlines())
    return f"<!doctype html><meta charset='utf-8'><body>{body}</body>"


def _write_csv(artifact: dict[str, Any], target: Path) -> None:
    payload = artifact.get("payload") if isinstance(artifact, dict) else {}
    payload = payload if isinstance(payload, dict) else {}
    columns = [str(item) for item in payload.get("columns", [])]
    rows = payload.get("rows", [])
    with target.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        if columns:
            writer.writerow(columns)
        for row in rows if isinstance(rows, list) else []:
            writer.writerow(row if isinstance(row, list) else [row])


def _render_svg(artifact: dict[str, Any]) -> str:
    title = html.escape(str(artifact.get("title") or "Artifact"))
    payload = artifact.get("payload") if isinstance(artifact, dict) else {}
    blocks = payload.get("blocks", []) if isinstance(payload, dict) else []
    block_items = [item for item in blocks if isinstance(item, dict)]
    height = max(220, 150 + len(block_items) * 120)
    lines = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="960" height="{height}" '
        f'viewBox="0 0 960 {height}" role="img">',
        f"<title>{title}</title>",
        '<rect width="960" height="100%" fill="#f8fafc"/>',
        f'<text x="48" y="72" font-family="Arial, sans-serif" font-size="34" '
        f'font-weight="700" fill="#111827">{title}</text>',
    ]
    for index, block in enumerate(block_items):
        y = 120 + index * 120
        block_title = html.escape(str(block.get("title") or f"Block {index + 1}"))
        block_text = html.escape(str(block.get("text") or block.get("summary") or ""))
        lines.extend(
            [
                f'<rect x="48" y="{y}" width="864" height="86" rx="8" '
                'fill="#ffffff" stroke="#cbd5e1"/>',
                f'<text x="72" y="{y + 34}" font-family="Arial, sans-serif" '
                f'font-size="22" font-weight="700" fill="#0f172a">{block_title}</text>',
                f'<text x="72" y="{y + 64}" font-family="Arial, sans-serif" '
                f'font-size="16" fill="#334155">{block_text}</text>',
            ]
        )
    lines.append("</svg>")
    return "\n".join(lines)


def _write_pptx(artifact: dict[str, Any], target: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    title = str(artifact.get("title") or "Artifact")
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    subtitle = title_slide.placeholders[1]
    subtitle.text = str(artifact.get("type") or "")
    for slide in _artifact_slides(artifact):
        content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        content_slide.shapes.title.text = slide["title"]
        body = content_slide.placeholders[1].text_frame
        body.clear()
        for index, bullet in enumerate(slide["bullets"]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
    presentation.save(str(target))


def _artifact_slides(artifact: dict[str, Any]) -> list[dict[str, list[str] | str]]:
    payload = artifact.get("payload") if isinstance(artifact, dict) else {}
    if not isinstance(payload, dict):
        return [{"title": "Artifact", "bullets": [str(payload or "")]}]
    raw_slides = payload.get("slides")
    if not isinstance(raw_slides, list):
        raw_slides = [
            slide
            for section in payload.get("sections", [])
            if isinstance(section, dict)
            for slide in section.get("slides", [])
        ]
    slides = [_normalize_slide(item) for item in raw_slides if isinstance(item, dict)]
    return slides or [{"title": "Artifact", "bullets": _fallback_bullets(payload)}]


def _normalize_slide(item: dict[str, Any]) -> dict[str, list[str] | str]:
    title = str(item.get("title") or "Untitled slide")
    bullets = item.get("bullets") or item.get("content") or item.get("points") or []
    if not isinstance(bullets, list):
        bullets = [bullets]
    return {"title": title, "bullets": [str(bullet) for bullet in bullets if bullet]}


def _fallback_bullets(payload: dict[str, Any]) -> list[str]:
    bullets: list[str] = []
    for key, value in payload.items():
        if isinstance(value, (str, int, float)):
            bullets.append(f"{key}: {value}")
    return bullets or ["No slide content available."]
