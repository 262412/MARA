from __future__ import annotations

import logging
import os
import re
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from ktem.preview.docx import extract_docx_text
from ktem.preview.context import PreviewPurpose, preview_access_for_user
from ktem.preview.office import OfficePreviewConversionService
from ktem.preview.service import PreviewService

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:
    Presentation = None
    MSO_SHAPE_TYPE = None


logger = logging.getLogger(__name__)

OFFICE_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".doc", ".ppt", ".xls"}


def detect_office_extension(file_name: str, file_path: str) -> str:
    ext = os.path.splitext((file_name or file_path or ""))[1].lower()
    if ext in OFFICE_EXTENSIONS:
        return ext

    if file_path and os.path.isfile(file_path):
        try:
            if zipfile.is_zipfile(file_path):
                with zipfile.ZipFile(file_path) as zf:
                    names = set(zf.namelist())
                if "word/document.xml" in names:
                    return ".docx"
                if "ppt/presentation.xml" in names:
                    return ".pptx"
                if "xl/workbook.xml" in names:
                    return ".xlsx"
        except Exception:
            pass

        try:
            with open(file_path, "rb") as file_obj:
                header = file_obj.read(8)
            if header.startswith(b"\xd0\xcf\x11\xe0"):
                return ".doc"
        except Exception:
            pass

    return ""


def read_text_file(file_path: str, max_chars: int = 9000) -> str:
    if not file_path or not os.path.isfile(file_path):
        return ""
    for enc in ("utf-8", "utf-16", "latin-1", "gbk"):
        try:
            with open(file_path, "r", encoding=enc, errors="ignore") as file_obj:
                content = file_obj.read(max_chars * 2)
            return content[:max_chars]
        except Exception:
            continue
    return ""


def extract_xlsx_text(file_path: str, max_chars: int = 9000) -> str:
    try:
        with zipfile.ZipFile(file_path) as zf:
            shared_strings: list[str] = []
            if "xl/sharedStrings.xml" in zf.namelist():
                with zf.open("xl/sharedStrings.xml") as file_obj:
                    ss_root = ET.fromstring(file_obj.read())
                for node in ss_root.iter():
                    if node.tag.endswith("}t") and node.text:
                        shared_strings.append(node.text)

            cells: list[str] = []
            total_chars = 0
            sheet_names = sorted(
                [
                    name
                    for name in zf.namelist()
                    if re.match(r"xl/worksheets/sheet\d+\.xml", name)
                ]
            )
            for sheet in sheet_names:
                with zf.open(sheet) as file_obj:
                    root = ET.fromstring(file_obj.read())
                for cell in root.iter():
                    if not cell.tag.endswith("}c"):
                        continue
                    cell_type = cell.attrib.get("t", "")
                    value = ""
                    for child in cell:
                        if child.tag.endswith("}v") and child.text:
                            value = child.text
                            break
                    if not value:
                        continue
                    if cell_type == "s":
                        try:
                            idx = int(value)
                            if 0 <= idx < len(shared_strings):
                                value = shared_strings[idx]
                        except Exception:
                            pass
                    cells.append(value)
                    total_chars += len(value)
                    if total_chars >= max_chars:
                        break
                if total_chars >= max_chars:
                    break
    except Exception:
        return ""
    return " ".join(cells)[:max_chars]


class PresentationTextService:
    def extract_slide_text(
        self, file_path: str, page: int, max_chars: int = 7000
    ) -> str:
        if not file_path or not os.path.isfile(file_path) or Presentation is None:
            return ""

        try:
            presentation = Presentation(file_path)
        except Exception:
            return ""

        slides = list(presentation.slides)
        if not slides:
            return ""
        page_idx = max(0, min(len(slides) - 1, int(page or 1) - 1))
        texts: list[str] = []
        total_chars = 0
        for shape in slides[page_idx].shapes:
            self._collect_shape_text(shape, texts)
            total_chars = sum(len(item) for item in texts)
            if total_chars >= max_chars:
                break
        return " ".join(part for part in texts if part).strip()[:max_chars]

    def _collect_shape_text(self, shape, texts: list[str]) -> None:
        if getattr(shape, "has_text_frame", False):
            for paragraph in getattr(shape.text_frame, "paragraphs", []):
                paragraph_text = (getattr(paragraph, "text", "") or "").strip()
                if paragraph_text:
                    texts.append(" ".join(paragraph_text.split()))

        if getattr(shape, "has_table", False):
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = (getattr(cell, "text", "") or "").strip()
                        if cell_text:
                            texts.append(" ".join(cell_text.split()))
            except Exception:
                pass

        try:
            shape_type = getattr(shape, "shape_type", None)
            if MSO_SHAPE_TYPE is not None and shape_type == MSO_SHAPE_TYPE.GROUP:
                for item in getattr(shape, "shapes", []):
                    self._collect_shape_text(item, texts)
        except Exception:
            pass


class PreviewFileResolver:
    def __init__(self, app, file_name_cache: dict[str, str]):
        self._app = app
        self._file_name_cache = file_name_cache
        self._service = PreviewService(app)

    @staticmethod
    def extract_first_selected_file_id(selected_file_ids):
        if not selected_file_ids:
            return ""

        selected = selected_file_ids[0]
        if isinstance(selected, str) and selected.startswith("["):
            try:
                import json

                selected_items = json.loads(selected)
                return selected_items[0] if selected_items else ""
            except Exception:
                return ""

        return selected

    def _access(self, user_id=None):
        return preview_access_for_user(self._app, user_id)

    def resolve_source(self, file_id: str, *, user_id=None):
        source = self._service.resolve_source(file_id, access=self._access(user_id))
        self._file_name_cache[file_id] = source.name
        return source

    def resolve_sources(self, file_ids, *, user_id=None, strict: bool = True):
        sources = self._service.resolve_sources(
            file_ids, access=self._access(user_id), strict=strict
        )
        self._file_name_cache.update(
            {source.file_id: source.name for source in sources}
        )
        return sources

    def resolve_file_path_by_id(self, file_id: str, *, user_id=None) -> str:
        if not file_id:
            return ""
        return str(self.resolve_source(file_id, user_id=user_id).path)

    def resolve_file_name_by_id(self, file_id: str, *, user_id=None) -> str:
        if not file_id:
            return ""
        return self.resolve_source(file_id, user_id=user_id).name

    def resolve_selected_file(self, selected_file_ids, *, user_id=None):
        file_id = self.extract_first_selected_file_id(selected_file_ids)
        if not file_id:
            return "", "", ""
        source = self.resolve_source(str(file_id), user_id=user_id)
        return source.file_id, source.name, str(source.path)


class PreviewSupportService:
    def __init__(self, app):
        self._app = app
        self._file_name_cache: dict[str, str] = {}
        self._resolver = PreviewFileResolver(app, self._file_name_cache)
        self._office_conversion = OfficePreviewConversionService(logger)
        self._presentation_service = PresentationTextService()

    def resolve_selected_file(
        self, selected_file_ids: list[str] | None, *, user_id=None
    ) -> tuple[str, str, str]:
        return self._resolver.resolve_selected_file(
            selected_file_ids or [], user_id=user_id
        )

    def resolve_file_path(self, file_id: str, *, user_id=None) -> str:
        return self._resolver.resolve_file_path_by_id(file_id, user_id=user_id)

    def resolve_file_name(self, file_id: str, *, user_id=None) -> str:
        return self._resolver.resolve_file_name_by_id(file_id, user_id=user_id)

    def resolve_sources(self, file_ids, *, user_id=None, strict: bool = True):
        return self._resolver.resolve_sources(file_ids, user_id=user_id, strict=strict)

    @staticmethod
    def extract_pdf_page_text(
        pdf_path: str, page_number: int, max_chars: int = 7000
    ) -> str:
        if not pdf_path or not os.path.isfile(pdf_path):
            return ""
        from pypdf import PdfReader

        try:
            reader = PdfReader(pdf_path)
            if not reader.pages:
                return ""
            page_idx = max(0, min(len(reader.pages) - 1, int(page_number or 1) - 1))
            text = reader.pages[page_idx].extract_text() or ""
            text = " ".join(str(text).split())
            return text[:max_chars]
        except Exception:
            return ""

    def get_page_context_text(
        self,
        file_id: str,
        file_name: str,
        page_number: int,
        max_chars: int = 7000,
        *,
        user_id=None,
    ) -> str:
        if not file_id or not file_name:
            return ""

        source = self._resolver.resolve_source(file_id, user_id=user_id)
        source_path = str(source.path)
        file_name = source.name

        source_extension = detect_office_extension(file_name, source_path)
        file_extension = (Path(file_name).suffix or Path(source_path).suffix).lower()

        fallback_text = ""
        if source_extension in {".pptx", ".ppt"}:
            fallback_text = self._presentation_service.extract_slide_text(
                source_path,
                page_number,
                max_chars=max_chars,
            )
        elif file_extension in {".docx", ".doc"}:
            fallback_text = extract_docx_text(source_path, max_chars=max_chars)
        elif file_extension in {".xlsx", ".xls", ".csv"}:
            fallback_text = extract_xlsx_text(source_path, max_chars=max_chars)
        elif file_extension in {".txt", ".md", ".html", ".mhtml"}:
            fallback_text = read_text_file(source_path, max_chars=max_chars)

        context = self._resolver._service.page_context(
            file_id,
            page=page_number,
            access=self._resolver._access(user_id),
            purpose=PreviewPurpose.DOCQA,
            max_chars=max_chars,
            fallback_text=fallback_text,
        )
        return context.text
