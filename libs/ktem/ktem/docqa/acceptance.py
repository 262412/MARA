from __future__ import annotations

import argparse
import contextlib
import json
import logging
import os
import shutil
import subprocess
import sys
import tempfile
import uuid
import warnings
from dataclasses import dataclass
from io import StringIO
from pathlib import Path
from typing import Any, cast


def _configure_noise_controls() -> None:
    os.environ.setdefault("GRPC_VERBOSITY", "ERROR")
    os.environ.setdefault("GLOG_minloglevel", "3")
    os.environ.setdefault("TF_CPP_MIN_LOG_LEVEL", "3")
    warnings.filterwarnings("ignore", message=".*ARC4 has been moved.*")
    warnings.filterwarnings("ignore", message=".*pkg_resources is deprecated.*")
    logging.getLogger("absl").setLevel(logging.ERROR)
    logging.getLogger("grpc").setLevel(logging.ERROR)


def _tail_text(text: str, max_lines: int = 40) -> str:
    lines = [line for line in text.splitlines() if line.strip()]
    if len(lines) <= max_lines:
        return "\n".join(lines)
    return "\n".join(lines[-max_lines:])


_configure_noise_controls()

from docx import Document as DocxDocument  # noqa: E402
from ktem.db.models import Conversation, engine  # noqa: E402
from ktem.docqa import DocQARequest, DocQARuntime  # noqa: E402
from ktem.index.file.index import FileIndex  # noqa: E402
from ktem.main import App  # noqa: E402
from ktem.utils.dependencies import DependencyChecker, find_soffice_binary  # noqa: E402
from openpyxl import Workbook  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from sqlmodel import Session, select  # noqa: E402


class AcceptanceFailure(RuntimeError):
    pass


@dataclass
class SampleFile:
    kind: str
    path: Path
    marker: str
    prompt: str


class AcceptanceMatrix:
    def __init__(self, keep_artifacts: bool = False, verbose: bool = False):
        self.keep_artifacts = keep_artifacts
        self.verbose = verbose
        self.completed_successfully = False
        self.runtime = DocQARuntime()
        self.user_id = str(self.runtime.user_id or "")
        self.work_dir = (
            Path(tempfile.gettempdir()) / f"kotaemon-acceptance-{uuid.uuid4().hex[:8]}"
        )
        self.samples_dir = self.work_dir / "samples"
        self.platform_dir = self.work_dir / "platform"
        self.graph_context_path = self.work_dir / "graph_context.json"
        self.sample_files: list[SampleFile] = []
        self.file_records: dict[str, dict[str, Any]] = {}
        self.cli_conversations: dict[str, str] = {}
        self.chat_conversation_id: str = ""
        self.web_conversation_id: str = ""
        self.created_conversation_ids: set[str] = set()
        self.created_file_ids: set[str] = set()
        self.results: list[dict[str, Any]] = []
        self.cli_executable = self._resolve_cli_executable()
        self.env = os.environ.copy()
        self.env["PYTHONIOENCODING"] = "utf-8"
        self.env.setdefault("GRPC_VERBOSITY", "ERROR")
        self.env.setdefault("GLOG_minloglevel", "3")
        self.env.setdefault("TF_CPP_MIN_LOG_LEVEL", "3")

    def _resolve_cli_executable(self) -> list[str]:
        return [sys.executable, "-m", "kotaemon.cli"]

    def _run_command(
        self,
        args: list[str],
        *,
        input_text: str | None = None,
        timeout: int = 900,
    ) -> subprocess.CompletedProcess[str]:
        completed = subprocess.run(
            args,
            env=self.env,
            input=input_text,
            text=True,
            capture_output=True,
            encoding="utf-8",
            timeout=timeout,
        )
        if completed.returncode != 0:
            raise AcceptanceFailure(
                "Command failed.\n"
                f"Command: {' '.join(args)}\n"
                f"STDOUT:\n{completed.stdout}\n"
                f"STDERR:\n{completed.stderr}"
            )
        return completed

    def _run_cli(
        self,
        *args: str,
        expect_json: bool = False,
        input_text: str | None = None,
        timeout: int = 900,
    ) -> Any:
        command = [*self.cli_executable, *args]
        completed = self._run_command(command, input_text=input_text, timeout=timeout)
        if not expect_json:
            return completed.stdout
        return self._extract_json_payload(completed.stdout)

    @staticmethod
    def _extract_json_payload(stdout: str) -> Any:
        lines = [line for line in stdout.splitlines() if line.strip()]
        errors: list[str] = []
        for index, line in enumerate(lines):
            stripped = line.lstrip()
            if not (stripped.startswith("{") or stripped.startswith("[")):
                continue
            payload = "\n".join(lines[index:])
            try:
                return json.loads(payload)
            except json.JSONDecodeError as exc:
                errors.append(f"line {index + 1}: {exc}")
        raise AcceptanceFailure(
            "Unable to parse JSON payload from CLI output.\n"
            f"Errors: {errors}\n"
            f"Raw stdout:\n{stdout}"
        )

    def _record(self, name: str, **details: Any) -> None:
        entry = {"name": name}
        entry.update(details)
        self.results.append(entry)

    def _assert_contains_marker(
        self, payload: dict[str, Any], marker: str, *, context: str
    ) -> None:
        joined = "\n".join(
            [
                str(payload.get("answer", "")),
                str(payload.get("references_text", "")),
                *[str(item) for item in payload.get("retrieval_messages", [])],
            ]
        )
        if marker not in joined:
            raise AcceptanceFailure(
                f"{context} did not contain expected marker '{marker}'.\n"
                f"Payload: {json.dumps(payload, ensure_ascii=False, indent=2)}"
            )

    def prepare_sample_files(self) -> None:
        self.samples_dir.mkdir(parents=True, exist_ok=True)
        self.platform_dir.mkdir(parents=True, exist_ok=True)

        graph_context = {
            "acceptance_run_id": self.work_dir.name,
            "related_markers": [
                "TXT-MARKER",
                "PDF-MARKER",
                "DOCX-MARKER",
                "XLSX-MARKER",
                "PPTX-MARKER",
            ],
        }
        self.graph_context_path.write_text(
            json.dumps(graph_context, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

        txt_path = self.samples_dir / "matrix_notes.txt"
        txt_marker = "TXT-MARKER-AURORA-2026"
        txt_path.write_text(
            "Acceptance sample for plain text.\n"
            f"Launch note marker: {txt_marker}.\n"
            "The review owner is Riley Stone.\n",
            encoding="utf-8",
        )
        self.sample_files.append(
            SampleFile(
                kind="txt",
                path=txt_path,
                marker=txt_marker,
                prompt="What is the launch note marker in the text file?",
            )
        )

        docx_path = self.samples_dir / "matrix_brief.docx"
        docx_marker = "DOCX-MARKER-ORBIT-2040"
        doc = DocxDocument()
        doc.add_heading("Acceptance Brief", level=1)
        doc.add_paragraph(f"Document marker: {docx_marker}.")
        doc.add_paragraph("Decision owner: Morgan Lee.")
        doc.save(docx_path)
        self.sample_files.append(
            SampleFile(
                kind="docx",
                path=docx_path,
                marker=docx_marker,
                prompt="What is the document marker in the DOCX file?",
            )
        )

        xlsx_path = self.samples_dir / "matrix_metrics.xlsx"
        xlsx_marker = "XLSX-MARKER-CIRRUS-77"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Metrics"
        sheet["A1"] = "Metric"
        sheet["B1"] = "Value"
        sheet["A2"] = "Workbook Marker"
        sheet["B2"] = xlsx_marker
        sheet["A3"] = "Forecast"
        sheet["B3"] = "June 2026"
        workbook.save(xlsx_path)
        self.sample_files.append(
            SampleFile(
                kind="xlsx",
                path=xlsx_path,
                marker=xlsx_marker,
                prompt="What is the workbook marker in the spreadsheet?",
            )
        )

        pptx_path = self.samples_dir / "matrix_story.pptx"
        pptx_marker = "PPTX-MARKER-EMBER-19"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Acceptance Deck"
        slide.placeholders[1].text = f"Slide marker: {pptx_marker}\nOwner: Jamie Frost"
        note_box = slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(6), Inches(1.5)
        )
        note_box.text_frame.text = "Deck validation content for CLI/Web acceptance."
        presentation.save(pptx_path)
        self.sample_files.append(
            SampleFile(
                kind="pptx",
                path=pptx_path,
                marker=pptx_marker,
                prompt="What is the slide marker in the presentation?",
            )
        )

        pdf_source_docx = self.samples_dir / "matrix_pdf_source.docx"
        pdf_marker = "PDF-MARKER-DELTA-55"
        pdf_doc = DocxDocument()
        pdf_doc.add_heading("Acceptance PDF Source", level=1)
        pdf_doc.add_paragraph(f"PDF marker: {pdf_marker}.")
        pdf_doc.add_paragraph("Release window: July 2026.")
        pdf_doc.save(pdf_source_docx)
        pdf_path = self.samples_dir / "matrix_report.pdf"
        self._convert_to_pdf(pdf_source_docx, pdf_path)
        self.sample_files.append(
            SampleFile(
                kind="pdf",
                path=pdf_path,
                marker=pdf_marker,
                prompt="What is the PDF marker in the report?",
            )
        )

        self._record(
            "prepare_samples",
            work_dir=str(self.work_dir),
            files=[str(sample.path) for sample in self.sample_files],
        )

    def _convert_to_pdf(self, source_path: Path, output_pdf: Path) -> None:
        available, libreoffice_info = DependencyChecker.check_libreoffice()
        if not available:
            raise AcceptanceFailure(
                "LibreOffice is required to generate the PDF sample."
            )

        soffice = find_soffice_binary()
        if not soffice:
            raise AcceptanceFailure(
                f"Unable to locate LibreOffice executable: {libreoffice_info}"
            )

        output_pdf.parent.mkdir(parents=True, exist_ok=True)
        command = [
            str(soffice),
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_pdf.parent),
            str(source_path),
        ]
        self._run_command(command, timeout=180)
        converted_pdf = output_pdf.parent / f"{source_path.stem}.pdf"
        if not converted_pdf.exists():
            raise AcceptanceFailure(
                f"LibreOffice did not produce PDF for {source_path}"
            )
        if converted_pdf != output_pdf:
            shutil.move(str(converted_pdf), str(output_pdf))

    def run_doctor(self) -> None:
        result = self._run_cli(
            "docqa", "doctor", "--json", expect_json=True, timeout=300
        )
        if not result.get("ok"):
            raise AcceptanceFailure(
                "DocQA doctor failed: "
                f"{json.dumps(result, ensure_ascii=False, indent=2)}"
            )
        self._record(
            "doctor",
            default_user_id=result.get("default_user_id"),
            index_name=result.get("index_name"),
            file_count=result.get("file_count"),
        )

    def run_index_and_file_listing(self) -> None:
        for sample in self.sample_files:
            result = self._run_cli(
                "docqa",
                "index",
                str(sample.path),
                "--json",
                expect_json=True,
                timeout=600,
            )
            if not result.get("successes"):
                raise AcceptanceFailure(
                    f"Indexing failed for {sample.path.name}: "
                    f"{json.dumps(result, ensure_ascii=False, indent=2)}"
                )
            self._record(
                "index",
                file=sample.path.name,
                successes=len(result.get("successes", [])),
            )

        files_payload = self._run_cli(
            "docqa", "files", "--json", expect_json=True, timeout=300
        )
        for sample in self.sample_files:
            matching = [
                row for row in files_payload if row.get("name") == sample.path.name
            ]
            if not matching:
                raise AcceptanceFailure(
                    f"Indexed file not found in file list: {sample.path.name}"
                )
            record = matching[0]
            self.file_records[sample.kind] = record
            self.created_file_ids.add(str(record["file_id"]))
        self._record(
            "files",
            indexed_files={
                kind: record["file_id"] for kind, record in self.file_records.items()
            },
        )

    def run_cli_ask_matrix(self) -> None:
        for sample in self.sample_files:
            record = self.file_records[sample.kind]
            payload = self._run_cli(
                "docqa",
                "ask",
                "--prompt",
                sample.prompt,
                "--file",
                str(record["file_id"]),
                "--active-file",
                str(record["file_id"]),
                "--page",
                "1",
                "--selected-text",
                sample.marker,
                "--graph-context-file",
                str(self.graph_context_path),
                "--citation",
                "inline",
                "--language",
                "en",
                "--json",
                expect_json=True,
                timeout=600,
            )
            conversation_id = str(payload.get("conversation_id") or "")
            if not conversation_id:
                raise AcceptanceFailure(
                    f"Missing conversation id for {sample.kind} ask payload."
                )
            self.created_conversation_ids.add(conversation_id)
            self.cli_conversations[sample.kind] = conversation_id
            self._assert_contains_marker(
                payload, sample.marker, context=f"CLI ask for {sample.kind}"
            )
            if (
                payload.get("graph_context", {}).get("acceptance_run_id")
                != self.work_dir.name
            ):
                raise AcceptanceFailure(
                    f"Graph context was not preserved for {sample.kind}."
                )
            self._record(
                "ask",
                file_kind=sample.kind,
                conversation_id=conversation_id,
                active_file_name=payload.get("active_file_name"),
            )

    def run_cli_chat_and_resume(self) -> None:
        txt_record = self.file_records["txt"]
        txt_sample = next(
            sample for sample in self.sample_files if sample.kind == "txt"
        )
        chat_output = self._run_cli(
            "docqa",
            "chat",
            "--file",
            str(txt_record["file_id"]),
            "--active-file",
            str(txt_record["file_id"]),
            "--page",
            "1",
            input_text=f"{txt_sample.prompt}\n/history\n/exit\n",
            timeout=600,
        )
        if txt_sample.marker not in chat_output:
            raise AcceptanceFailure("docqa chat output did not include the TXT marker.")

        match = None
        for line in chat_output.splitlines():
            if line.startswith("Conversation:"):
                match = line.split(":", 1)[1].strip()
                break
        if not match:
            raise AcceptanceFailure(
                "docqa chat output did not include a conversation id."
            )
        self.chat_conversation_id = match
        self.created_conversation_ids.add(match)
        self._record("chat", conversation_id=match)

        pdf_sample = next(
            sample for sample in self.sample_files if sample.kind == "pdf"
        )
        pdf_conversation_id = self.cli_conversations["pdf"]
        resume_output = self._run_cli(
            "docqa",
            "resume",
            pdf_conversation_id,
            input_text="/history\n/exit\n",
            timeout=300,
        )
        if pdf_sample.marker not in resume_output:
            raise AcceptanceFailure(
                "docqa resume output did not include the PDF marker."
            )
        self._record("resume", conversation_id=pdf_conversation_id)

    def run_cli_to_web_restore_matrix(self) -> None:
        app = App()
        app.make()
        chat_page = cast(Any, app).chat_page
        file_index = next(
            index for index in app.index_manager.indices if isinstance(index, FileIndex)
        )
        selector_ui = file_index.get_selector_component_ui()
        _, selector_choices = selector_ui.load_files([], self.user_id)

        history = chat_page.chat_control.load_chat_history(self.user_id)
        history_ids = {option[1] for option in history}

        preview_runtime = self.runtime._preview
        for sample in self.sample_files:
            conversation_id = self.cli_conversations[sample.kind]
            file_id = str(self.file_records[sample.kind]["file_id"])
            if conversation_id not in history_ids:
                raise AcceptanceFailure(
                    f"Web history does not include CLI conversation {conversation_id}."
                )

            selected = chat_page.chat_control.select_conv(conversation_id, self.user_id)
            restored_conversation_id = selected[0]
            restored_messages = selected[3]

            offset = 11
            restored_selector_map: dict[int, Any] = {}
            for index in app.index_manager.indices:
                if index.selector is None:
                    continue
                if isinstance(index.selector, tuple):
                    width = len(index.default_selector)
                    restored_selector_map[index.id] = list(
                        selected[offset : offset + width]
                    )
                    offset += width
                else:
                    restored_selector_map[index.id] = selected[offset]
                    offset += 1

            file_selector_state = restored_selector_map[file_index.id]
            restored_selected_ids = file_index.resolve_selected_ids(
                self.user_id, file_selector_state
            )
            graph_source_ids = chat_page.load_conversation_graph_state(conversation_id)
            rows, list_html, focus_label = chat_page.refresh_chat_file_list(
                conversation_id,
                self.user_id,
                selector_choices,
                restored_selected_ids,
                graph_source_ids,
                "",
            )
            preview = chat_page.page_preview.refresh_selected_file_preview(
                selector_choices,
                restored_selected_ids,
                1,
                1,
            )
            (
                preview_file_id,
                preview_file_name,
                preview_file_path,
                preview_page,
                _preview_total_pages,
                preview_src,
                preview_notice,
            ) = preview
            page_context = preview_runtime.get_page_context_text(
                file_id, sample.path.name, 1
            )

            if restored_conversation_id != conversation_id:
                raise AcceptanceFailure(
                    f"Web restored wrong conversation for {sample.kind}."
                )
            if not restored_messages:
                raise AcceptanceFailure(
                    f"Web restored empty message history for {sample.kind}."
                )
            if sample.marker not in restored_messages[-1][1]:
                raise AcceptanceFailure(
                    f"Web restored answer missing marker for {sample.kind}."
                )
            if file_id not in restored_selected_ids:
                raise AcceptanceFailure(
                    f"Web restored selected ids missing file for {sample.kind}."
                )
            if file_id not in graph_source_ids:
                raise AcceptanceFailure(
                    f"Web restored graph ids missing file for {sample.kind}."
                )
            if not any(row.get("id") == file_id for row in rows):
                raise AcceptanceFailure(
                    f"Web chat file list missing file for {sample.kind}."
                )
            if sample.path.name not in focus_label:
                raise AcceptanceFailure(
                    f"Web focus label missing file name for {sample.kind}."
                )
            if preview_file_id != file_id:
                raise AcceptanceFailure(
                    f"Web preview restored wrong file id for {sample.kind}."
                )
            if preview_file_name != sample.path.name:
                raise AcceptanceFailure(
                    f"Web preview restored wrong file name for {sample.kind}."
                )
            if not preview_file_path or not Path(preview_file_path).exists():
                raise AcceptanceFailure(
                    f"Web preview resolved invalid file path for {sample.kind}."
                )
            if int(preview_page or 0) != 1:
                raise AcceptanceFailure(
                    f"Web preview restored wrong page for {sample.kind}."
                )
            if not preview_src and not preview_notice:
                raise AcceptanceFailure(
                    "Web preview returned neither content nor notice "
                    f"for {sample.kind}."
                )
            if sample.marker not in (page_context or ""):
                raise AcceptanceFailure(
                    f"Preview page context missing marker for {sample.kind}."
                )
            if sample.path.name not in list_html:
                raise AcceptanceFailure(
                    f"Web list HTML missing file name for {sample.kind}."
                )

            self._record(
                "cli_to_web",
                file_kind=sample.kind,
                conversation_id=conversation_id,
                preview_has_content=bool(preview_src),
                preview_notice=preview_notice,
            )

    def run_web_to_cli_resume(self) -> None:
        app = App()
        app.make()
        chat_page = cast(Any, app).chat_page

        sample = next(item for item in self.sample_files if item.kind == "docx")
        file_record = self.file_records["docx"]
        conversation_id, _dropdown_update = chat_page.chat_control.new_conv(
            self.user_id
        )
        if not conversation_id:
            raise AcceptanceFailure("Web conversation creation failed.")
        self.web_conversation_id = str(conversation_id)
        self.created_conversation_ids.add(self.web_conversation_id)

        response = chat_page.docqa.run_turn(
            DocQARequest(
                prompt=sample.prompt,
                conversation_id=self.web_conversation_id,
                selected_file_ids=[str(file_record["file_id"])],
                active_file_id=str(file_record["file_id"]),
                active_file_name=sample.path.name,
                page_number=1,
                selected_text=sample.marker,
                graph_context={"created_via": "web_acceptance"},
                user_id=self.user_id,
                origin="web",
            )
        )
        self._assert_contains_marker(
            response.as_dict(), sample.marker, context="Web-origin conversation"
        )

        resume_output = self._run_cli(
            "docqa",
            "resume",
            self.web_conversation_id,
            input_text="/history\n/exit\n",
            timeout=300,
        )
        if sample.marker not in resume_output:
            raise AcceptanceFailure(
                "CLI resume could not reopen the Web-created conversation."
            )

        with Session(engine) as session:
            row = session.exec(
                select(Conversation).where(Conversation.id == self.web_conversation_id)
            ).one()
            data_source = dict(row.data_source or {})
            if data_source.get("origin") != "web":
                raise AcceptanceFailure(
                    "Web-created conversation lost its web origin marker."
                )
        self._record("web_to_cli", conversation_id=self.web_conversation_id)

    def run_platform_matrix(self) -> None:
        self._run_cli("platform", "validate", timeout=300)
        self._record("platform_bundle_validate", status="pass")

        platform_cases: list[tuple[str, str, list[str]]] = [
            ("codex", "minimal", []),
            ("codex", "full", []),
            ("codex", "selective", ["skills", "agents", "AGENTS.md"]),
            ("claude-code", "minimal", []),
            ("claude-code", "full", []),
            ("claude-code", "selective", ["skills", "agents", "CLAUDE.md", "commands"]),
        ]

        for platform_name, mode, items in platform_cases:
            target_dir = self.platform_dir / f"{platform_name}-{mode}"
            command = [
                "platform",
                "install",
                "--platform",
                platform_name,
                "--mode",
                mode,
                "--target-dir",
                str(target_dir),
                "--yes",
            ]
            for item in items:
                command.extend(["--item", item])
            install_output = self._run_cli(*command, timeout=300)

            validate_output = self._run_cli(
                "platform",
                "validate",
                "--platform",
                platform_name,
                "--installed",
                "--target-dir",
                str(target_dir),
                timeout=300,
            )
            expected_skill = target_dir / "skills" / "kotaemon-docqa" / "SKILL.md"
            if not expected_skill.exists():
                raise AcceptanceFailure(
                    f"Installed platform bundle missing docqa skill: {expected_skill}"
                )
            if platform_name == "claude-code":
                expected_command = target_dir / "commands" / "kotaemon-docqa.md"
                if mode in {"full", "selective"} and not expected_command.exists():
                    raise AcceptanceFailure(
                        "Installed Claude Code bundle missing docqa command "
                        f"wrapper: {expected_command}"
                    )
            if f"{platform_name}: PASS" not in validate_output:
                raise AcceptanceFailure(
                    "Installed platform validation did not pass for "
                    f"{platform_name} {mode}.\n{validate_output}"
                )
            self._record(
                "platform_install",
                platform=platform_name,
                mode=mode,
                target_dir=str(target_dir),
                install_output=install_output.strip().splitlines()[:5],
            )

    def run_delete_and_cleanup_validation(self) -> None:
        delete_payload = self._run_cli(
            "docqa",
            "delete",
            *sorted(self.created_file_ids),
            "--json",
            expect_json=True,
            timeout=300,
        )
        deleted_ids = {str(item.get("file_id") or "") for item in delete_payload}
        missing = sorted(self.created_file_ids - deleted_ids)
        if missing:
            raise AcceptanceFailure(
                f"docqa delete did not remove expected file ids: {missing}"
            )

        remaining_files = self._run_cli(
            "docqa", "files", "--json", expect_json=True, timeout=300
        )
        remaining_ids = {str(item.get("file_id") or "") for item in remaining_files}
        still_present = sorted(self.created_file_ids & remaining_ids)
        if still_present:
            raise AcceptanceFailure(
                f"Deleted file ids still present after delete: {still_present}"
            )
        self._record("delete", deleted_count=len(delete_payload))

    def cleanup(self) -> None:
        for conversation_id in sorted(self.created_conversation_ids):
            with Session(engine) as session:
                row = session.exec(
                    select(Conversation).where(Conversation.id == conversation_id)
                ).one_or_none()
                if row is not None:
                    session.delete(row)
                    session.commit()

        if (
            not self.keep_artifacts
            and self.completed_successfully
            and self.work_dir.exists()
        ):
            shutil.rmtree(self.work_dir, ignore_errors=True)

    def run(self) -> dict[str, Any]:
        try:
            self.prepare_sample_files()
            self.run_doctor()
            self.run_index_and_file_listing()
            self.run_cli_ask_matrix()
            self.run_cli_chat_and_resume()
            self.run_cli_to_web_restore_matrix()
            self.run_web_to_cli_resume()
            self.run_platform_matrix()
            self.run_delete_and_cleanup_validation()
            self.completed_successfully = True
            return {
                "status": "pass",
                "user_id": self.user_id,
                "work_dir": str(self.work_dir) if self.keep_artifacts else "",
                "results": self.results,
            }
        finally:
            self.cleanup()


def main() -> int:
    parser = argparse.ArgumentParser(
        description="Run the DocQA cross-entry acceptance matrix."
    )
    parser.add_argument(
        "--keep-artifacts",
        action="store_true",
        help="Keep temporary documents and install targets after the run.",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="Show in-process logs and warnings instead of suppressing them.",
    )
    args = parser.parse_args()

    matrix = AcceptanceMatrix(
        keep_artifacts=args.keep_artifacts,
        verbose=args.verbose,
    )
    stdout_buffer = StringIO()
    stderr_buffer = StringIO()
    try:
        if matrix.verbose:
            result = matrix.run()
        else:
            with (
                contextlib.redirect_stdout(stdout_buffer),
                contextlib.redirect_stderr(stderr_buffer),
            ):
                result = matrix.run()
    except Exception as exc:
        failure = {
            "status": "fail",
            "error": str(exc),
            "partial_results": matrix.results,
            "work_dir": str(matrix.work_dir),
        }
        captured_stdout = stdout_buffer.getvalue()
        captured_stderr = stderr_buffer.getvalue()
        if captured_stdout.strip():
            failure["captured_stdout_tail"] = _tail_text(captured_stdout)
        if captured_stderr.strip():
            failure["captured_stderr_tail"] = _tail_text(captured_stderr)
        print(json.dumps(failure, ensure_ascii=False, indent=2))
        return 1

    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
