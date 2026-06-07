import sys
from types import ModuleType

from ktem.docqa.artifact_exports import (
    export_artifact_to_path,
    render_artifact_markdown,
)


def test_render_artifact_markdown_formats_payload_and_citations():
    text = render_artifact_markdown(
        {
            "artifact_id": "artifact-1",
            "type": "study_guide",
            "title": "Study guide",
            "payload": {"overview": "Grounded overview"},
            "citations": [{"citation_id": "c1", "source_id": "file-1"}],
        }
    )

    assert "# Study guide" in text
    assert "## Overview" in text
    assert "Grounded overview" in text
    assert "## Citations" in text


def test_export_artifact_to_markdown_writes_grounded_metadata(tmp_path):
    output = tmp_path / "guide.md"

    result = export_artifact_to_path(
        {
            "artifact_id": "artifact-1",
            "type": "study_guide",
            "title": "Study guide",
            "payload": {"overview": "Grounded overview"},
            "citations": [{"citation_id": "c1", "source_id": "file-1"}],
        },
        export_format="md",
        output_path=output,
    )

    assert result == output
    text = output.read_text(encoding="utf-8")
    assert "# Study guide" in text
    assert "Grounded overview" in text
    assert "c1" in text


def test_export_data_table_to_csv(tmp_path):
    output = tmp_path / "table.csv"

    export_artifact_to_path(
        {
            "artifact_id": "artifact-2",
            "type": "data_table",
            "title": "Table",
            "payload": {
                "columns": ["Metric", "Value"],
                "rows": [["Revenue", "$10"], ["Margin", "20%"]],
            },
        },
        export_format="csv",
        output_path=output,
    )

    assert output.read_text(encoding="utf-8").splitlines() == [
        "Metric,Value",
        "Revenue,$10",
        "Margin,20%",
    ]


def test_export_infographic_to_svg(tmp_path):
    output = tmp_path / "infographic.svg"

    export_artifact_to_path(
        {
            "artifact_id": "artifact-3",
            "type": "infographic",
            "title": "Revenue infographic",
            "payload": {
                "blocks": [
                    {
                        "title": "Growth",
                        "text": "Revenue increased after source-grounded retrieval.",
                    }
                ]
            },
        },
        export_format="svg",
        output_path=output,
    )

    text = output.read_text(encoding="utf-8")
    assert text.startswith("<svg")
    assert "Revenue infographic" in text
    assert "Growth" in text
    assert "source-grounded retrieval" in text


def test_export_slide_deck_to_pptx(tmp_path):
    from pptx import Presentation

    output = tmp_path / "artifact-deck.pptx"

    export_artifact_to_path(
        {
            "artifact_id": "artifact-4",
            "type": "slide_deck",
            "title": "MARA Evidence Deck",
            "payload": {
                "slides": [
                    {
                        "title": "Grounded Retrieval",
                        "bullets": ["Evidence stays linked to sources."],
                    }
                ]
            },
        },
        export_format="pptx",
        output_path=output,
    )

    presentation = Presentation(str(output))
    slide_text = "\n".join(
        shape.text for slide in presentation.slides for shape in slide.shapes
    )
    assert "MARA Evidence Deck" in slide_text
    assert "Grounded Retrieval" in slide_text
    assert "Evidence stays linked to sources." in slide_text


def test_export_media_requires_adapter(tmp_path):
    output = tmp_path / "overview.mp3"

    try:
        export_artifact_to_path(
            {
                "artifact_id": "artifact-5",
                "type": "audio_overview",
                "title": "Audio",
                "payload": {"media_status": "script_only", "script": []},
            },
            export_format="mp3",
            output_path=output,
        )
    except ValueError as exc:
        assert "requires a configured media export adapter" in str(exc)
    else:
        raise AssertionError("mp3 export should require a configured adapter")


def test_export_media_uses_configured_adapter(tmp_path):
    output = tmp_path / "overview.mp3"
    calls = []

    def media_adapter(artifact, export_format, output_path):
        calls.append((artifact["artifact_id"], export_format, output_path))
        output_path.write_bytes(b"audio-bytes")
        return output_path

    result = export_artifact_to_path(
        {
            "artifact_id": "artifact-6",
            "type": "audio_overview",
            "title": "Audio",
            "payload": {"media_status": "script_only", "script": []},
        },
        export_format="mp3",
        output_path=output,
        media_export_adapter=media_adapter,
    )

    assert result == output
    assert output.read_bytes() == b"audio-bytes"
    assert calls == [("artifact-6", "mp3", output)]


def test_export_media_uses_configured_dotted_adapter(monkeypatch, tmp_path):
    output = tmp_path / "overview.mp3"
    calls = []

    def media_adapter(artifact, export_format, output_path):
        calls.append((artifact["artifact_id"], export_format, output_path))
        output_path.write_bytes(b"configured-audio")
        return output_path

    monkeypatch.setattr(
        "ktem.docqa.artifact_exports.configured_media_export_adapter",
        lambda: media_adapter,
    )

    result = export_artifact_to_path(
        {
            "artifact_id": "artifact-7",
            "type": "audio_overview",
            "title": "Audio",
            "payload": {"media_status": "script_only", "script": []},
        },
        export_format="mp3",
        output_path=output,
    )

    assert result == output
    assert output.read_bytes() == b"configured-audio"
    assert calls == [("artifact-7", "mp3", output)]


def test_configured_media_export_adapter_imports_dotted_setting(monkeypatch):
    from ktem.docqa import artifact_exports
    from theflow.settings import settings as flowsettings

    module = ModuleType("fake_media_export_adapter")

    def media_adapter(_artifact, _export_format, _output_path):
        return _output_path

    setattr(module, "media_adapter", media_adapter)
    monkeypatch.setitem(sys.modules, "fake_media_export_adapter", module)
    monkeypatch.setattr(
        flowsettings,
        "KH_MARA_ARTIFACT_MEDIA_EXPORT_ADAPTER",
        "fake_media_export_adapter.media_adapter",
        raising=False,
    )

    assert artifact_exports.configured_media_export_adapter() is media_adapter
