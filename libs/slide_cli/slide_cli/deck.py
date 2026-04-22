from __future__ import annotations

from dataclasses import dataclass, field
import os
from pathlib import Path
import shutil
import subprocess
from typing import Callable

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover - guarded at runtime
    Presentation = None
    MSO_SHAPE_TYPE = None


TITLE_PLACEHOLDER_NAMES = {"TITLE", "CENTER_TITLE"}
BODY_PLACEHOLDER_NAMES = {"SUBTITLE", "BODY", "OBJECT"}


PathLike = str | Path


@dataclass(frozen=True, slots=True)
class ShapeSnapshot:
    slide_number: int
    target_id: str
    text: str
    kind: str
    shape_name: str = ""
    shape_id: int | None = None
    placeholder_type: str = ""
    parent_target_id: str = ""

    def summary_line(self, max_chars: int = 120) -> str:
        preview = _clip_text(self.text, max_chars=max_chars)
        return f"{self.target_id} [{self.kind}] {preview}"

    def as_dict(self) -> dict[str, object]:
        return {
            "slide_number": self.slide_number,
            "target_id": self.target_id,
            "text": self.text,
            "kind": self.kind,
            "shape_name": self.shape_name,
            "shape_id": self.shape_id,
            "placeholder_type": self.placeholder_type,
            "parent_target_id": self.parent_target_id,
        }


@dataclass(frozen=True, slots=True)
class SlideSnapshot:
    slide_number: int
    title: str
    shapes: tuple[ShapeSnapshot, ...]

    def summary_text(self, max_chars: int = 120) -> str:
        lines = [f"Slide {self.slide_number}: {self.title or '(untitled)'}"]
        lines.extend(shape.summary_line(max_chars=max_chars) for shape in self.shapes)
        return "\n".join(lines)

    def as_dict(self) -> dict[str, object]:
        return {
            "slide_number": self.slide_number,
            "title": self.title,
            "shapes": [shape.as_dict() for shape in self.shapes],
        }


@dataclass(frozen=True, slots=True)
class DeckSnapshot:
    source_path: Path
    slides: tuple[SlideSnapshot, ...]

    @property
    def slide_count(self) -> int:
        return len(self.slides)

    def summary_text(self, max_chars: int = 120) -> str:
        return "\n\n".join(
            slide.summary_text(max_chars=max_chars) for slide in self.slides
        )

    def as_dict(self) -> dict[str, object]:
        return {
            "source_path": str(self.source_path),
            "slide_count": self.slide_count,
            "slides": [slide.as_dict() for slide in self.slides],
        }


@dataclass(slots=True)
class TextReplaceOp:
    slide_number: int
    target_id: str
    before_text: str | None = None
    after_text: str = ""

    def as_dict(self) -> dict[str, object]:
        return {
            "slide_number": self.slide_number,
            "target_id": self.target_id,
            "before_text": self.before_text,
            "after_text": self.after_text,
        }


@dataclass(slots=True)
class DeckPatch:
    summary: str
    edits: list[TextReplaceOp] = field(default_factory=list)

    def as_dict(self) -> dict[str, object]:
        return {
            "summary": self.summary,
            "edits": [edit.as_dict() for edit in self.edits],
        }


@dataclass(slots=True)
class DeckEditResult:
    output_path: Path
    written: bool
    applied_target_ids: list[str] = field(default_factory=list)
    skipped: list[str] = field(default_factory=list)

    @property
    def applied_count(self) -> int:
        return len(self.applied_target_ids)

    @property
    def skipped_count(self) -> int:
        return len(self.skipped)

    def as_dict(self) -> dict[str, object]:
        return {
            "output_path": str(self.output_path),
            "written": self.written,
            "applied_target_ids": list(self.applied_target_ids),
            "applied_count": self.applied_count,
            "skipped": list(self.skipped),
            "skipped_count": self.skipped_count,
        }


@dataclass(slots=True)
class _TextTargetBinding:
    slide_number: int
    target_id: str
    get_text: Callable[[], str]
    set_text: Callable[[str], None]


def load_deck_snapshot(source_path: PathLike) -> DeckSnapshot:
    path = Path(source_path)
    presentation = _open_presentation(path)
    slides = tuple(
        _snapshot_slide(slide_number, slide)
        for slide_number, slide in enumerate(presentation.slides, start=1)
    )
    return DeckSnapshot(source_path=path, slides=slides)


def apply_deck_patch(
    source_path: PathLike,
    patch: DeckPatch,
    *,
    output_path: PathLike | None = None,
) -> DeckEditResult:
    path = Path(source_path)
    presentation = _open_presentation(path)
    bindings = _collect_bindings(presentation)
    applied: list[str] = []
    skipped: list[str] = []

    for edit in patch.edits:
        binding = bindings.get(edit.target_id)
        if binding is None:
            skipped.append(f"{edit.target_id}: target not found")
            continue
        if binding.slide_number != edit.slide_number:
            skipped.append(
                f"{edit.target_id}: belongs to slide {binding.slide_number}, not {edit.slide_number}"
            )
            continue

        current_text = binding.get_text()
        if edit.before_text is not None and current_text != edit.before_text:
            skipped.append(f"{edit.target_id}: before_text mismatch")
            continue

        binding.set_text(edit.after_text)
        applied.append(edit.target_id)

    destination = Path(output_path) if output_path is not None else _default_output_path(path)
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(destination))

    return DeckEditResult(
        output_path=destination,
        written=True,
        applied_target_ids=applied,
        skipped=skipped,
    )


def export_deck_pdf(
    source_path: PathLike,
    *,
    output_path: PathLike | None = None,
    soffice_path: str | None = None,
    timeout_sec: int = 120,
) -> Path:
    source = Path(source_path).resolve()
    if not source.exists():
        raise FileNotFoundError(source)

    resolved_soffice = soffice_path or os.environ.get("SOFFICE_PATH") or shutil.which("soffice")
    if not resolved_soffice:
        raise RuntimeError("LibreOffice is required to export slide decks to PDF.")

    requested_output = Path(output_path).resolve() if output_path is not None else None
    target_dir = requested_output.parent if requested_output is not None else source.parent
    target_dir.mkdir(parents=True, exist_ok=True)

    completed = subprocess.run(
        [
            str(resolved_soffice),
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(target_dir),
            str(source),
        ],
        capture_output=True,
        text=True,
        encoding="utf-8",
        timeout=timeout_sec,
        check=False,
    )
    if completed.returncode != 0:
        details = completed.stderr.strip() or completed.stdout.strip() or "unknown error"
        raise RuntimeError(f"LibreOffice export failed: {details}")

    converted_path = target_dir / f"{source.stem}.pdf"
    if not converted_path.exists():
        raise RuntimeError("LibreOffice export did not create the expected PDF output.")

    if requested_output is not None and converted_path != requested_output:
        requested_output.parent.mkdir(parents=True, exist_ok=True)
        converted_path.replace(requested_output)
        return requested_output
    return converted_path


def _open_presentation(path: Path):
    if Presentation is None:  # pragma: no cover - depends on optional install
        raise RuntimeError("python-pptx is required to load slide decks")
    if not path.exists():
        raise FileNotFoundError(path)
    return Presentation(str(path))


def _snapshot_slide(slide_number: int, slide) -> SlideSnapshot:
    shapes = tuple(_iter_text_targets(slide_number, slide))
    return SlideSnapshot(
        slide_number=slide_number,
        title=_detect_slide_title(slide, shapes),
        shapes=shapes,
    )


def _iter_text_targets(slide_number: int, slide) -> list[ShapeSnapshot]:
    snapshots: list[ShapeSnapshot] = []
    for index, shape in enumerate(slide.shapes, start=1):
        snapshots.extend(
            _snapshot_shape_targets(
                slide_number=slide_number,
                shape=shape,
                shape_path=(_shape_key(shape, index),),
                parent_target_id="",
            )
        )
    return snapshots


def _snapshot_shape_targets(
    *,
    slide_number: int,
    shape,
    shape_path: tuple[str, ...],
    parent_target_id: str,
) -> list[ShapeSnapshot]:
    snapshots: list[ShapeSnapshot] = []

    if _is_group_shape(shape):
        for index, child in enumerate(getattr(shape, "shapes", []), start=1):
            snapshots.extend(
                _snapshot_shape_targets(
                    slide_number=slide_number,
                    shape=child,
                    shape_path=shape_path + (_shape_key(child, index),),
                    parent_target_id=parent_target_id,
                )
            )
        return snapshots

    shape_id = _shape_id(shape)
    shape_name = _shape_name(shape)
    placeholder_type = _placeholder_type_name(shape)

    if getattr(shape, "has_table", False):
        table_target_id = _target_id(slide_number, shape_path, "table")
        try:
            for row_index, row in enumerate(shape.table.rows, start=1):
                for column_index, cell in enumerate(row.cells, start=1):
                    snapshots.append(
                        ShapeSnapshot(
                            slide_number=slide_number,
                            target_id=_target_id(
                                slide_number,
                                shape_path,
                                f"table-r{row_index}c{column_index}",
                            ),
                            text=getattr(cell, "text", "") or "",
                            kind="table_cell",
                            shape_name=f"{shape_name}[r{row_index}c{column_index}]",
                            shape_id=shape_id,
                            placeholder_type=placeholder_type,
                            parent_target_id=table_target_id,
                        )
                    )
        except Exception:
            return snapshots
        return snapshots

    if getattr(shape, "has_text_frame", False):
        snapshots.append(
            ShapeSnapshot(
                slide_number=slide_number,
                target_id=_target_id(slide_number, shape_path, "text"),
                text=getattr(shape, "text", "") or "",
                kind=_shape_kind(shape),
                shape_name=shape_name,
                shape_id=shape_id,
                placeholder_type=placeholder_type,
                parent_target_id=parent_target_id,
            )
        )

    return snapshots


def _collect_bindings(presentation) -> dict[str, _TextTargetBinding]:
    bindings: dict[str, _TextTargetBinding] = {}
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for index, shape in enumerate(slide.shapes, start=1):
            for binding in _collect_shape_bindings(
                slide_number=slide_number,
                shape=shape,
                shape_path=(_shape_key(shape, index),),
            ):
                bindings[binding.target_id] = binding
    return bindings


def _collect_shape_bindings(
    *,
    slide_number: int,
    shape,
    shape_path: tuple[str, ...],
) -> list[_TextTargetBinding]:
    bindings: list[_TextTargetBinding] = []

    if _is_group_shape(shape):
        for index, child in enumerate(getattr(shape, "shapes", []), start=1):
            bindings.extend(
                _collect_shape_bindings(
                    slide_number=slide_number,
                    shape=child,
                    shape_path=shape_path + (_shape_key(child, index),),
                )
            )
        return bindings

    if getattr(shape, "has_table", False):
        try:
            for row_index, row in enumerate(shape.table.rows, start=1):
                for column_index, cell in enumerate(row.cells, start=1):
                    bindings.append(
                        _TextTargetBinding(
                            slide_number=slide_number,
                            target_id=_target_id(
                                slide_number,
                                shape_path,
                                f"table-r{row_index}c{column_index}",
                            ),
                            get_text=lambda cell=cell: getattr(cell, "text", "") or "",
                            set_text=lambda value, cell=cell: setattr(
                                cell,
                                "text",
                                value or "",
                            ),
                        )
                    )
        except Exception:
            return bindings
        return bindings

    if getattr(shape, "has_text_frame", False):
        bindings.append(
            _TextTargetBinding(
                slide_number=slide_number,
                target_id=_target_id(slide_number, shape_path, "text"),
                get_text=lambda shape=shape: getattr(shape, "text", "") or "",
                set_text=lambda value, shape=shape: setattr(
                    shape,
                    "text",
                    value or "",
                ),
            )
        )

    return bindings


def _detect_slide_title(slide, shapes: tuple[ShapeSnapshot, ...]) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        title_text = getattr(title_shape, "text", "") or ""
        if title_text.strip():
            return title_text

    for shape in shapes:
        if shape.kind == "title" and shape.text.strip():
            return shape.text

    return ""


def _shape_kind(shape) -> str:
    placeholder_type = _placeholder_type_name(shape)
    if placeholder_type in TITLE_PLACEHOLDER_NAMES:
        return "title"
    if placeholder_type in BODY_PLACEHOLDER_NAMES:
        return "body"
    return "text"


def _shape_key(shape, index: int) -> str:
    shape_id = _shape_id(shape)
    if shape_id is not None:
        return f"shape-{shape_id}"
    return f"idx-{index}"


def _shape_id(shape) -> int | None:
    shape_id = getattr(shape, "shape_id", None)
    if shape_id is None:
        return None
    try:
        return int(shape_id)
    except (TypeError, ValueError):
        return None


def _shape_name(shape) -> str:
    return (getattr(shape, "name", "") or "").strip()


def _placeholder_type_name(shape) -> str:
    try:
        if not getattr(shape, "is_placeholder", False):
            return ""
        placeholder_type = getattr(shape.placeholder_format, "type", None)
    except Exception:
        return ""
    if placeholder_type is None:
        return ""
    return getattr(placeholder_type, "name", str(placeholder_type))


def _is_group_shape(shape) -> bool:
    if MSO_SHAPE_TYPE is None:
        return False
    try:
        return getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _target_id(slide_number: int, shape_path: tuple[str, ...], leaf: str) -> str:
    return f"slide-{slide_number}/{'/'.join(shape_path)}/{leaf}"


def _default_output_path(source_path: Path) -> Path:
    return source_path.with_name(f"{source_path.stem}.patched{source_path.suffix}")


def _clip_text(text: str, *, max_chars: int) -> str:
    compact = " ".join((text or "").split())
    if not compact:
        return "(empty)"
    if len(compact) <= max_chars:
        return compact
    return f"{compact[: max_chars - 3]}..."
