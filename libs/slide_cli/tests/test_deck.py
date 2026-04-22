from pptx import Presentation

from slide_cli.deck import DeckPatch, TextReplaceOp, apply_deck_patch, load_deck_snapshot


def test_load_snapshot_and_apply_patch(tmp_path):
    source_path = tmp_path / "source.pptx"
    output_path = tmp_path / "rewritten.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Business Review"
    slide.placeholders[1].text = "Revenue is flat.\nPipeline quality is mixed."
    presentation.save(source_path)

    snapshot = load_deck_snapshot(source_path)

    assert snapshot.slide_count == 1
    assert snapshot.slides[0].title == "Quarterly Business Review"
    assert len(snapshot.slides[0].shapes) >= 2

    title_shape = next(shape for shape in snapshot.slides[0].shapes if shape.text)
    patch = DeckPatch(
        summary="Sharpen executive framing",
        edits=[
            TextReplaceOp(
                slide_number=1,
                target_id=title_shape.target_id,
                before_text=title_shape.text,
                after_text="Executive Quarterly Business Review",
            )
        ],
    )

    result = apply_deck_patch(source_path, patch, output_path=output_path)
    rewritten = Presentation(output_path)

    assert result.written is True
    assert result.output_path == output_path
    assert rewritten.slides[0].shapes.title.text == "Executive Quarterly Business Review"
