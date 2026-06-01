import json
from types import SimpleNamespace

from pptx import Presentation
from slide_cli.agent import SlideAgentRunner


def test_agent_runner_can_use_tool_before_final(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "QBR"
    slide.placeholders[1].text = "Revenue is flat."
    presentation.save(deck_path)

    responses = [
        SimpleNamespace(
            text=(
                "Thought: I should inspect the first slide before rewriting it.\n"
                "Action: read_slide\n"
                "Action Input: 1"
            )
        ),
        SimpleNamespace(
            text=(
                "Thought: I now know the final answer.\n"
                "Final Answer: "
                + json.dumps(
                    {
                        "assistant_response": "Reframed the title for executives.",
                        "patch": {
                            "summary": "Update slide title",
                            "edits": [
                                {
                                    "slide_number": 1,
                                    "target_id": "slide-1/shape-2/text",
                                    "before_text": "QBR",
                                    "after_text": "Executive QBR",
                                }
                            ],
                        },
                    }
                )
            )
        ),
    ]

    monkeypatch.setattr(
        "slide_cli.agent.run_completion",
        lambda **kwargs: responses.pop(0),
    )

    runner = SlideAgentRunner(
        input_path=str(deck_path),
        model="gpt-4o-mini",
        config_path="missing.yml",
        cwd=str(tmp_path),
    )
    result = runner.run("Rewrite the deck title for executives.")

    assert result["assistant_response"] == "Reframed the title for executives."
    assert result["patch"].summary == "Update slide title"
    assert result["patch"].edits[0].after_text == "Executive QBR"
    assert result["observations"]


def test_agent_runner_supports_phase_one_tools(monkeypatch, tmp_path):
    from slide_cli.config import SlideAgentConfig

    deck_path = tmp_path / "deck.pptx"
    pdf_path = tmp_path / "deck.pdf"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "QBR"
    slide.placeholders[1].text = "Revenue is flat."
    presentation.save(deck_path)

    monkeypatch.setattr(
        "slide_cli.agent.export_deck_pdf",
        lambda source_path, output_path=None, **kwargs: pdf_path,
    )

    runner = SlideAgentRunner(
        input_path=str(deck_path),
        config=SlideAgentConfig(
            cwd=str(tmp_path),
            config_path="missing.yml",
        ),
    )

    extracted = runner._execute_tool("extract_slide_text", "1")
    review = runner._execute_tool("review_deck", "")
    write_result = runner._execute_tool(
        "write_file",
        json.dumps({"path": "notes.txt", "content": "Hello from mara-research-cli"}),
    )
    export_result = runner._execute_tool(
        "export_pdf",
        json.dumps({"output_path": str(pdf_path)}),
    )

    assert "Revenue is flat." in extracted
    assert "slide_count" in review
    assert (tmp_path / "notes.txt").read_text(
        encoding="utf-8"
    ) == "Hello from mara-research-cli"
    assert "notes.txt" in write_result
    assert str(pdf_path) in export_result


def test_agent_runner_instruction_uses_top_level_agent_line_language(
    monkeypatch, tmp_path
):
    deck_path = tmp_path / "deck.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "QBR"
    slide.placeholders[1].text = "Revenue is flat."
    presentation.save(deck_path)

    monkeypatch.setattr(
        "slide_cli.agent.run_completion",
        lambda **kwargs: SimpleNamespace(
            text=(
                "Thought: I now know the final answer.\n"
                "Final Answer: "
                + json.dumps(
                    {
                        "assistant_response": "No deck patch required.",
                        "patch": {"summary": "No deck edits", "edits": []},
                    }
                )
            )
        ),
    )

    runner = SlideAgentRunner(
        input_path=str(deck_path),
        model="gpt-4o-mini",
        config_path="missing.yml",
        cwd=str(tmp_path),
    )
    instruction = runner._build_instruction(
        user_prompt="Review the workspace notes before deciding on deck changes.",
        history_text="(none)",
    )

    assert "top-level MARA CLI agent line" in instruction
    assert "high-permission workflow" in instruction
    assert "workspace-side file changes" in instruction
