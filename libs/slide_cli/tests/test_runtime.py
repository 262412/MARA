from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from slide_cli.deck import DeckPatch, TextReplaceOp
from slide_cli.session_store import SlideSessionStore


def _make_deck(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "QBR"
    slide.placeholders[1].text = "Revenue is flat."
    presentation.save(path)


def test_run_slide_task_preview_mode_does_not_write(monkeypatch, tmp_path):
    from slide_cli import runtime

    deck_path = tmp_path / "deck.pptx"
    _make_deck(deck_path)

    monkeypatch.setattr(
        runtime,
        "SlideSessionStore",
        lambda *args, **kwargs: SlideSessionStore(base_dir=tmp_path / "runtime"),
    )

    class StubRunner:
        def __init__(self, **kwargs):
            self.kwargs = kwargs

        def run(self, prompt, history=None):
            return {
                "assistant_response": "Prepared a rewrite preview.",
                "patch": DeckPatch(
                    summary="Rewrite title",
                    edits=[
                        TextReplaceOp(
                            slide_number=1,
                            target_id="slide-1/shape-2/text",
                            before_text="QBR",
                            after_text="Executive QBR",
                        )
                    ],
                ),
                "observations": [],
                "raw_responses": [],
            }

    monkeypatch.setattr(runtime, "SlideAgentRunner", StubRunner)

    result = runtime.run_slide_task(
        input_path=str(deck_path),
        prompt="Rewrite this for executives",
        apply_mode="preview",
    )

    assert result["mode"] == "preview"
    assert result["output_path"] == ""
    assert result["can_apply"] is True
    assert not deck_path.with_name("deck.rewritten.pptx").exists()


def test_apply_session_patch_writes_latest_patch(monkeypatch, tmp_path):
    from slide_cli import runtime

    deck_path = tmp_path / "deck.pptx"
    output_path = tmp_path / "deck.final.pptx"
    _make_deck(deck_path)

    monkeypatch.setattr(
        runtime,
        "SlideSessionStore",
        lambda *args, **kwargs: SlideSessionStore(base_dir=tmp_path / "runtime"),
    )

    class StubRunner:
        def __init__(self, **kwargs):
            self.kwargs = kwargs

        def run(self, prompt, history=None):
            return {
                "assistant_response": "Prepared a rewrite preview.",
                "patch": DeckPatch(
                    summary="Rewrite title",
                    edits=[
                        TextReplaceOp(
                            slide_number=1,
                            target_id="slide-1/shape-2/text",
                            before_text="QBR",
                            after_text="Executive QBR",
                        )
                    ],
                ),
                "observations": [],
                "raw_responses": [],
            }

    monkeypatch.setattr(runtime, "SlideAgentRunner", StubRunner)

    preview = runtime.run_slide_task(
        input_path=str(deck_path),
        prompt="Rewrite this for executives",
        apply_mode="confirm",
    )

    applied = runtime.apply_session_patch(
        preview["session_id"],
        output_path=str(output_path),
        base_dir=tmp_path / "runtime",
    )

    rewritten = Presentation(output_path)

    assert applied["output_path"] == str(output_path)
    assert applied["applied_count"] == 1
    assert rewritten.slides[0].shapes.title.text == "Executive QBR"
