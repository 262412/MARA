import json
import sys

import pytest
from click.testing import CliRunner
from pptx import Presentation
from slide_cli.cli import main


def _listed_commands(help_output: str) -> set[str]:
    commands: set[str] = set()
    in_commands = False
    for raw_line in help_output.splitlines():
        line = raw_line.rstrip()
        if line.startswith("Commands:"):
            in_commands = True
            continue
        if not in_commands:
            continue
        if not line.strip():
            break
        stripped = line.strip()
        if not stripped:
            continue
        commands.add(stripped.split()[0])
    return commands


def test_help_lists_core_commands():
    runner = CliRunner()

    result = runner.invoke(main, ["--help"], terminal_width=300)

    assert result.exit_code == 0, result.output
    commands = _listed_commands(result.output)
    for token in [
        "inspect",
        "read-slide",
        "extract",
        "search",
        "apply",
        "export-pdf",
        "review",
        "doctor",
        "files",
        "read",
        "write",
        "delete",
        "shell",
        "run",
        "chat",
        "sessions",
        "resume",
        "docqa",
    ]:
        assert token in commands
    for token in [
        "ask",
        "index",
        "docqa-sessions",
        "resume-docqa",
    ]:
        assert token not in commands


@pytest.mark.parametrize(
    ("command_name", "expected_tokens"),
    [
        ("inspect", ["--file", "--json"]),
        ("read-slide", ["--file", "--slide", "--json"]),
        ("extract", ["--file", "--slide", "--json"]),
        ("search", ["--file", "--query", "--json"]),
        ("apply", ["session_id", "--output", "--json"]),
        ("export-pdf", ["--file", "--output", "--json"]),
        ("review", ["--file", "--json"]),
    ],
)
def test_phase3_command_help_exposes_canonical_arguments(command_name, expected_tokens):
    runner = CliRunner()

    result = runner.invoke(main, [command_name, "--help"])

    assert result.exit_code == 0, result.output
    for token in expected_tokens:
        assert token in result.output


def test_help_describes_two_line_product_model():
    runner = CliRunner()

    result = runner.invoke(main, ["--help"], terminal_width=300)

    assert result.exit_code == 0, result.output
    for token in [
        "Unified slide product CLI.",
        "Top-level agent line:",
        "Specialist DocQA line:",
        "Support lines:",
        "`slide run`",
        "read-only deck observability commands",
        "`slide app ...`",
        "`slide model ...`",
        "`slide platform ...`",
        "`slide docqa ...`",
    ]:
        assert token in result.output


def test_run_help_describes_high_permission_agent_workflow():
    runner = CliRunner()

    result = runner.invoke(main, ["run", "--help"])

    assert result.exit_code == 0, result.output
    for token in [
        "Run one high-permission slide agent workflow.",
        "--file",
        "--prompt",
        "--approval-policy",
    ]:
        assert token in result.output


def test_chat_help_describes_interactive_agent_workflow():
    runner = CliRunner()

    result = runner.invoke(main, ["chat", "--help"])

    assert result.exit_code == 0, result.output
    for token in [
        "Open an interactive high-permission slide agent session.",
        "--file",
        "--prompt",
        "--approval-policy",
    ]:
        assert token in result.output


def test_workspace_commands_round_trip(tmp_path):
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    runner = CliRunner()

    write_result = runner.invoke(
        main,
        [
            "write",
            "notes.txt",
            "--content",
            "hello from slide",
            "--cwd",
            str(workspace),
        ],
    )
    files_result = runner.invoke(
        main,
        ["files", "--cwd", str(workspace), "--json"],
    )
    read_result = runner.invoke(
        main,
        ["read", "notes.txt", "--cwd", str(workspace)],
    )
    shell_result = runner.invoke(
        main,
        [
            "shell",
            "--command",
            f'"{sys.executable}" -c "print(321)"',
            "--cwd",
            str(workspace),
            "--json",
        ],
    )
    delete_result = runner.invoke(
        main,
        ["delete", "notes.txt", "--cwd", str(workspace), "--yes"],
    )

    assert write_result.exit_code == 0, write_result.output
    assert "Wrote" in write_result.output

    assert files_result.exit_code == 0, files_result.output
    files_payload = json.loads(files_result.output)
    assert "notes.txt" in files_payload["paths"]

    assert read_result.exit_code == 0, read_result.output
    assert "hello from slide" in read_result.output

    assert shell_result.exit_code == 0, shell_result.output
    shell_payload = json.loads(shell_result.output)
    assert shell_payload["returncode"] == 0
    assert "321" in shell_payload["stdout"]

    assert delete_result.exit_code == 0, delete_result.output
    assert "Deleted file notes.txt" in delete_result.output


def test_apply_json_applies_saved_session_patch(monkeypatch, tmp_path):
    captured: dict[str, object] = {}
    output_path = tmp_path / "rewritten.pptx"

    def _fake_apply_session_patch(session_id, *, output_path=None, base_dir=None):
        captured.update(
            session_id=session_id,
            output_path=output_path,
            base_dir=base_dir,
        )
        return {
            "session_id": session_id,
            "output_path": str(output_path or ""),
            "applied_count": 1,
        }

    monkeypatch.setattr("slide_cli.cli.apply_session_patch", _fake_apply_session_patch)

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "apply",
            "session-123",
            "--output",
            str(output_path),
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    assert captured["session_id"] == "session-123"
    assert captured["output_path"] == str(output_path)
    decoded = json.loads(result.output)
    assert decoded["session_id"] == "session-123"
    assert decoded["output_path"] == str(output_path)
    assert decoded["applied_count"] == 1


def test_inspect_json_emits_deck_summary(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    captured: dict[str, object] = {}

    def _fake_inspect_slide_deck(input_path):
        captured["input_path"] = input_path
        return {
            "input_path": input_path,
            "slide_count": 1,
            "summary": "Slide 1: Quarterly Business Review",
        }

    monkeypatch.setattr("slide_cli.cli.inspect_slide_deck", _fake_inspect_slide_deck)

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "inspect",
            "--file",
            str(deck_path),
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    assert captured["input_path"] == str(deck_path)
    decoded = json.loads(result.output)
    assert decoded["slide_count"] == 1
    assert "Quarterly Business Review" in decoded["summary"]


def test_read_slide_json_emits_slide_summary(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    captured: dict[str, object] = {}

    def _fake_read_slide_summary(input_path, *, slide_number):
        captured["input_path"] = input_path
        captured["slide_number"] = slide_number
        return {
            "input_path": input_path,
            "slide_number": slide_number,
            "summary": "Slide 2: Revenue outlook",
        }

    monkeypatch.setattr("slide_cli.cli.read_slide_summary", _fake_read_slide_summary)

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "read-slide",
            "--file",
            str(deck_path),
            "--slide",
            "2",
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    assert captured["input_path"] == str(deck_path)
    assert captured["slide_number"] == 2
    decoded = json.loads(result.output)
    assert decoded["slide_number"] == 2
    assert "Revenue outlook" in decoded["summary"]


def test_extract_json_emits_plain_text(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    captured: dict[str, object] = {}

    def _fake_extract_slide_text(input_path, *, slide_number=None):
        captured["input_path"] = input_path
        captured["slide_number"] = slide_number
        return {
            "input_path": input_path,
            "slide_number": slide_number,
            "text": "Slide 1: Quarterly Business Review",
        }

    monkeypatch.setattr("slide_cli.cli.extract_slide_text", _fake_extract_slide_text)

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "extract",
            "--file",
            str(deck_path),
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    assert captured["input_path"] == str(deck_path)
    assert captured["slide_number"] is None
    decoded = json.loads(result.output)
    assert "Quarterly Business Review" in decoded["text"]


def test_search_json_emits_match_list(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    captured: dict[str, object] = {}

    def _fake_search_slide_deck(input_path, *, query):
        captured["input_path"] = input_path
        captured["query"] = query
        return {
            "input_path": input_path,
            "query": query,
            "count": 1,
            "matches": ["Slide 1: Revenue grew 20%"],
        }

    monkeypatch.setattr("slide_cli.cli.search_slide_deck", _fake_search_slide_deck)

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "search",
            "--file",
            str(deck_path),
            "--query",
            "Revenue",
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    assert captured["input_path"] == str(deck_path)
    assert captured["query"] == "Revenue"
    decoded = json.loads(result.output)
    assert decoded["count"] == 1
    assert "Revenue grew 20%" in decoded["matches"][0]


def test_export_pdf_json_emits_output_path(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    output_path = tmp_path / "deck.pdf"
    captured: dict[str, object] = {}

    def _fake_export_deck_pdf(source_path, *, output_path=None, **kwargs):
        captured.update(
            source_path=source_path,
            output_path=output_path,
            kwargs=kwargs,
        )
        return output_path

    monkeypatch.setattr(
        "slide_cli.cli.export_deck_pdf", _fake_export_deck_pdf, raising=False
    )
    monkeypatch.setattr("slide_cli.deck.export_deck_pdf", _fake_export_deck_pdf)

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "export-pdf",
            "--file",
            str(deck_path),
            "--output",
            str(output_path),
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    assert captured["source_path"] == str(deck_path)
    assert captured["output_path"] == str(output_path)
    decoded = json.loads(result.output)
    assert decoded["output_path"] == str(output_path)


def test_review_json_emits_structured_deck_summary(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Quarterly Business Review"
    presentation.save(deck_path)
    captured: dict[str, object] = {}

    def _fake_review_deck(self):
        captured["file"] = getattr(self, "input_path", None)
        return json.dumps(
            {
                "slide_count": 1,
                "text_target_count": 2,
                "untitled_slides": [],
                "empty_targets": [],
                "long_targets": [],
            },
            ensure_ascii=False,
            indent=2,
        )

    monkeypatch.setattr(
        "slide_cli.tools.SlideToolContext.review_deck", _fake_review_deck
    )

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "review",
            "--file",
            str(deck_path),
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    decoded = json.loads(result.output)
    assert decoded["slide_count"] == 1
    assert decoded["text_target_count"] == 2


def test_top_level_docqa_aliases_are_not_available():
    runner = CliRunner()

    for command_name in ["ask", "index", "docqa-sessions", "resume-docqa"]:
        result = runner.invoke(main, [command_name, "--help"])
        assert result.exit_code != 0
        assert "No such command" in result.output


def test_top_level_docqa_help_includes_mainline_action_guide():
    runner = CliRunner()

    result = runner.invoke(main, ["docqa", "--help"])

    assert result.exit_code == 0, result.output
    for token in [
        "Inspect indexed files",
        "Delete indexed files",
        "Inspect saved sessions",
        "Maintainer acceptance check",
    ]:
        assert token in result.output


def test_doctor_json(monkeypatch):
    payload = {
        "ok": True,
        "config_path": "C:/tmp/modelcli.yml",
        "providers": {"openai": {"available": False, "reason": "missing API key"}},
        "python_pptx": True,
        "libreoffice": False,
    }
    monkeypatch.setattr("slide_cli.cli._collect_doctor_payload", lambda: payload)

    runner = CliRunner()
    result = runner.invoke(main, ["doctor", "--json"])

    assert result.exit_code == 0, result.output
    decoded = json.loads(result.output)
    assert decoded["ok"] is True
    assert decoded["python_pptx"] is True
    assert decoded["providers"]["openai"]["reason"] == "missing API key"


def test_run_dry_run_json(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")

    monkeypatch.setattr(
        "slide_cli.cli.run_slide_task",
        lambda **kwargs: {
            "status": "ok",
            "mode": "dry-run",
            "response": "Suggested edits prepared.",
            "output_path": "",
            "patch": {"summary": "Rewrite title", "edits": []},
            "session_id": "session-1",
            "input_path": kwargs["input_path"],
        },
    )

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "run",
            "--file",
            str(deck_path),
            "--prompt",
            "Rewrite this deck for executives",
            "--dry-run",
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    decoded = json.loads(result.output)
    assert decoded["status"] == "ok"
    assert decoded["mode"] == "dry-run"
    assert decoded["session_id"] == "session-1"
    assert decoded["input_path"] == str(deck_path)


def test_run_defaults_to_preview_mode(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    captured: dict[str, object] = {}

    def _fake_run_slide_task(**kwargs):
        captured.update(kwargs)
        return {
            "status": "ok",
            "mode": "preview",
            "response": "Preview ready.",
            "output_path": "",
            "patch": {"summary": "Rewrite title", "edits": []},
            "session_id": "session-2",
            "input_path": kwargs["input_path"],
            "can_apply": True,
        }

    monkeypatch.setattr("slide_cli.cli.run_slide_task", _fake_run_slide_task)

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "run",
            "--file",
            str(deck_path),
            "--prompt",
            "Rewrite this deck for executives",
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    assert captured["apply_mode"] == "preview"


def test_run_apply_flag_sets_apply_mode(monkeypatch, tmp_path):
    deck_path = tmp_path / "deck.pptx"
    deck_path.write_bytes(b"placeholder")
    captured: dict[str, object] = {}

    def _fake_run_slide_task(**kwargs):
        captured.update(kwargs)
        return {
            "status": "ok",
            "mode": "apply",
            "response": "Applied the patch.",
            "output_path": str(tmp_path / "rewritten.pptx"),
            "patch": {"summary": "Rewrite title", "edits": []},
            "session_id": "session-3",
            "input_path": kwargs["input_path"],
            "can_apply": False,
        }

    monkeypatch.setattr("slide_cli.cli.run_slide_task", _fake_run_slide_task)

    runner = CliRunner()
    result = runner.invoke(
        main,
        [
            "run",
            "--file",
            str(deck_path),
            "--prompt",
            "Rewrite this deck for executives",
            "--apply",
            "--json",
        ],
    )

    assert result.exit_code == 0, result.output
    assert captured["apply_mode"] == "apply"
