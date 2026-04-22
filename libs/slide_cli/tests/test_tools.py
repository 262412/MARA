import subprocess

from pptx import Presentation

from slide_cli.deck import load_deck_snapshot
from slide_cli.tools import ReadFileTool, RunShellTool, SlideToolContext


def _make_context(tmp_path):
    deck_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "QBR"
    slide.placeholders[1].text = "Revenue is flat."
    presentation.save(deck_path)
    return SlideToolContext(
        input_path=deck_path,
        workspace_root=tmp_path,
        snapshot=load_deck_snapshot(deck_path),
        shell_timeout_sec=1,
    )


def test_read_file_blocks_paths_outside_workspace(tmp_path):
    context = _make_context(tmp_path)
    outside_path = tmp_path.parent / "outside.txt"
    outside_path.write_text("secret", encoding="utf-8")

    tool = ReadFileTool(tool_ctx=context)
    result = tool.run(str(outside_path))

    assert "outside the workspace root" in result


def test_run_shell_reports_timeout(monkeypatch, tmp_path):
    context = _make_context(tmp_path)
    tool = RunShellTool(tool_ctx=context)

    def _raise_timeout(*args, **kwargs):
        raise subprocess.TimeoutExpired(cmd="python -c pass", timeout=1)

    monkeypatch.setattr("slide_cli.tools.subprocess.run", _raise_timeout)

    result = tool.run("python -c \"import time; time.sleep(5)\"")

    assert "timed out after 1 seconds" in result
