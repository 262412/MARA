from __future__ import annotations

import gc
import tempfile
import unittest
import zipfile
from pathlib import Path

from .application import DesktopApplicationService
from .smoke_fixture import (
    GATE2_SMOKE_FILE_ID,
    GATE2_SMOKE_SESSION_ID,
    GATE3_FORMAT_INPUT_NAMES,
    seed_smoke_fixture,
)


def assert_gate3_format_inputs(test: unittest.TestCase, data_root: Path) -> None:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation

    input_root = data_root / "tmp"
    test.assertTrue(
        all((input_root / name).is_file() for name in GATE3_FORMAT_INPUT_NAMES)
    )
    document = Document(input_root / "gate3-format.docx")
    test.assertIn(
        "MARA modern Office format matrix fixture.",
        [paragraph.text for paragraph in document.paragraphs],
    )
    workbook = load_workbook(input_root / "gate3-format.xlsx", read_only=True)
    try:
        test.assertEqual(workbook["MARA"]["A2"].value, "indexed_files")
    finally:
        workbook.close()
    presentation = Presentation(input_root / "gate3-format.pptx")
    test.assertEqual(presentation.slides[0].shapes.title.text, "Gate 3 PPTX")
    from kotaemon.indices.ingests.files import KH_DEFAULT_FILE_EXTRACTORS

    for name, marker in (
        ("gate3-format.docx", "MARA modern Office format matrix fixture."),
        ("gate3-format.xlsx", "indexed_files"),
        ("gate3-format.pptx", "MARA modern Office format matrix fixture."),
    ):
        reader = KH_DEFAULT_FILE_EXTRACTORS[Path(name).suffix]
        documents = reader.load_data(input_root / name)
        searchable_text = "\n".join(document.text for document in documents)
        test.assertIn(marker, searchable_text, name)
    with zipfile.ZipFile(input_root / "gate3-format.zip") as archive:
        test.assertEqual(archive.namelist(), ["gate3-zip-note.md"])


def assert_gate2_session_detail(test: unittest.TestCase, session: dict) -> None:
    test.assertEqual(
        session["messages"],
        [
            {"role": "user", "content": "What is this fixture?"},
            {
                "role": "assistant",
                "content": "A packaged Desktop smoke record.",
            },
        ],
    )
    test.assertNotIn("data_source", session)
    test.assertNotIn("user_id", session)


class Gate2SmokeFixtureTest(unittest.TestCase):
    def test_seeds_one_file_and_one_session_for_the_real_application_service(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as temporary_directory:
            data_root = Path(temporary_directory) / "MARA"

            seed_smoke_fixture(data_root)
            seed_smoke_fixture(data_root)

            from ktem.db.models import engine

            service: DesktopApplicationService | None = None
            try:
                self.assertFalse(
                    (
                        data_root
                        / "state"
                        / "ktem_app_data"
                        / "user_data"
                        / "vectorstore"
                        / "chroma.sqlite3"
                    ).exists()
                )
                service = DesktopApplicationService()
                doctor = service.get_doctor()
                files = service.list_files()
                sessions = service.list_sessions()
                session = service.get_session(GATE2_SMOKE_SESSION_ID)
                import_capabilities = service.get_import_capabilities()

                self.assertTrue(doctor["ok"])
                self.assertEqual(doctor["file_count"], 1)
                self.assertEqual(doctor["session_count"], 1)
                self.assertEqual(
                    [record["file_id"] for record in files],
                    [GATE2_SMOKE_FILE_ID],
                )
                self.assertNotIn("path", files[0])
                self.assertEqual(
                    [record["conversation_id"] for record in sessions],
                    [GATE2_SMOKE_SESSION_ID],
                )
                assert_gate2_session_detail(self, session)
                self.assertTrue(
                    {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".zip"}
                    <= set(import_capabilities["supported_extensions"])
                )
                assert_gate3_format_inputs(self, data_root)
                self.assertEqual(
                    service.delete_file(GATE2_SMOKE_FILE_ID),
                    [
                        {
                            "file_id": GATE2_SMOKE_FILE_ID,
                            "name": "gate2-smoke.txt",
                        }
                    ],
                )
                self.assertEqual(service.list_files(), [])
                self.assertFalse(
                    (
                        data_root
                        / "state"
                        / "ktem_app_data"
                        / "user_data"
                        / "files"
                        / "index_1"
                        / "gate2-smoke.txt"
                    ).exists()
                )
            finally:
                service = None
                from chromadb.api.client import SharedSystemClient

                for system in tuple(SharedSystemClient._identifier_to_system.values()):
                    system.stop()
                SharedSystemClient.clear_system_cache()
                engine.dispose()
                gc.collect()


if __name__ == "__main__":
    unittest.main()
