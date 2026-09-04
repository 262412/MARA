"""Generate hostile binary preview fixtures in a disposable directory."""

from __future__ import annotations

import argparse
import os
import shutil
import tempfile
import zipfile
from pathlib import Path

import fitz
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter

PDF_PAYLOAD = """</span><img src=x onerror="parent.__maraPdfXss+=1"><script>parent.__maraPdfXss+=10</script><svg onload="parent.__maraPdfXss+=100">"""
DOCX_PAYLOAD = """<img src=x onerror="parent.__maraDocxXss+=1"><script>parent.__maraDocxXss+=10</script>"""
PPTX_PAYLOAD = """</p><img src=x onerror="parent.__maraPptxXss+=1"><script>parent.__maraPptxXss+=10</script>"""
ANSWER_PAYLOAD = r"""<img src=x onerror="window.__maraAnswerXss+=1">
<script>window.__maraAnswerXss+=10</script>
[JS LINK](javascript:window.__maraAnswerXss+=100)
![SVG IMAGE](data:image/svg+xml,<svg onload="window.__maraAnswerXss+=1000"></svg>)
[DATA HTML](data:text/html,<script>window.__maraAnswerXss+=10000</script>)
</summary></details><form action="https://attacker.invalid/answer-form"><button autofocus onfocus="window.__maraAnswerXss+=100000">submit</button></form>
Math: $\href{javascript:window.__maraAnswerXss+=1000000}{click}$"""
SVG_PAYLOAD = b"""<svg xmlns="http://www.w3.org/2000/svg" onload="parent.__maraOfficeXss+=1000"><script>parent.__maraOfficeXss+=10000</script></svg>"""


def _add_docx_hyperlink(paragraph, label: str, target: str) -> None:
    relation_id = paragraph.part.relate_to(
        target,
        RELATIONSHIP_TYPE.HYPERLINK,
        is_external=True,
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relation_id)
    run = OxmlElement("w:r")
    text = OxmlElement("w:t")
    text.text = label
    run.append(text)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _append_svg_relationship(path: Path, relationship_part: str) -> None:
    file_descriptor, temporary_name = tempfile.mkstemp(suffix=path.suffix)
    os.close(file_descriptor)
    temporary = Path(temporary_name)
    try:
        with zipfile.ZipFile(path) as source, zipfile.ZipFile(
            temporary,
            "w",
            zipfile.ZIP_DEFLATED,
        ) as target:
            for item in source.infolist():
                payload = source.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    payload = payload.replace(
                        b"</Types>",
                        b'<Default Extension="svg" ContentType="image/svg+xml"/></Types>',
                    )
                if item.filename == relationship_part:
                    relationship_target = (
                        b"media/evil.svg"
                        if path.suffix == ".docx"
                        else b"../media/evil.svg"
                    )
                    payload = payload.replace(
                        b"</Relationships>",
                        b'<Relationship Id="rIdMaraSvg" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="'
                        + relationship_target
                        + b'"/></Relationships>',
                    )
                target.writestr(item, payload)
            media_root = "word" if path.suffix == ".docx" else "ppt"
            target.writestr(f"{media_root}/media/evil.svg", SVG_PAYLOAD)
        shutil.move(temporary, path)
    finally:
        temporary.unlink(missing_ok=True)


def build_pdf(path: Path) -> None:
    document = fitz.open()
    texts = [PDF_PAYLOAD, "MARA PDF PAGE TWO", "MARA PDF PAGE THREE"]
    for index, text in enumerate(texts):
        page = document.new_page()
        page.insert_textbox(fitz.Rect(36, 50, 560, 760), text, fontsize=8)
        if index == 0:
            page.insert_link(
                {
                    "kind": fitz.LINK_URI,
                    "from": fitz.Rect(36, 200, 260, 225),
                    "uri": "javascript:parent.__maraPdfXss+=1000",
                }
            )
            page.insert_link(
                {
                    "kind": fitz.LINK_URI,
                    "from": fitz.Rect(36, 230, 260, 255),
                    "uri": "https://attacker.invalid/pdf-link",
                }
            )
    raw_path = path.with_suffix(".raw.pdf")
    document.save(raw_path)
    document.close()
    reader = PdfReader(raw_path)
    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    writer.add_js("app.launchURL('https://attacker.invalid/pdf-open', true)")
    with path.open("wb") as output:
        writer.write(output)
    raw_path.unlink()


def build_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("MARA DOCX SAFE TEXT")
    document.add_paragraph(DOCX_PAYLOAD)
    links = document.add_paragraph()
    _add_docx_hyperlink(
        links,
        "DOCX JS LINK",
        "javascript:parent.__maraDocxXss+=100",
    )
    _add_docx_hyperlink(
        links,
        "DOCX HTTP LINK",
        "https://attacker.invalid/docx-link",
    )
    document.core_properties.comments = (
        '"><img src=x onerror="parent.__maraDocxXss+=100000">'
    )
    document.save(path)
    _append_svg_relationship(path, "word/_rels/document.xml.rels")


def build_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(4))
    box.text_frame.text = PPTX_PAYLOAD
    paragraph = box.text_frame.add_paragraph()
    js_run = paragraph.add_run()
    js_run.text = "PPTX JS LINK"
    js_run.hyperlink.address = "javascript:parent.__maraPptxXss+=100"
    http_run = paragraph.add_run()
    http_run.text = " PPTX HTTP LINK"
    http_run.hyperlink.address = "https://attacker.invalid/pptx-link"
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(8), Inches(2)
    ).text_frame.text = "MARA PPTX SLIDE TWO"
    presentation.save(path)
    _append_svg_relationship(path, "ppt/slides/_rels/slide1.xml.rels")


def build_fixtures(output: Path) -> None:
    output.mkdir(parents=True, exist_ok=True)
    build_pdf(output / "malicious.pdf")
    build_docx(output / "malicious.docx")
    build_pptx(output / "malicious.pptx")
    marker = '<img src=x onerror="window.__maraNoticeXss=1">'
    (output / f"corrupt-{marker}.docx").write_bytes(b"not-a-docx")
    (output / f"corrupt-{marker}.pptx").write_bytes(b"not-a-pptx")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()
    build_fixtures(args.output)


if __name__ == "__main__":
    main()
