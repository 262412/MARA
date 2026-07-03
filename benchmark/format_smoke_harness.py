from __future__ import annotations

import csv
import json
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Callable
from zipfile import BadZipFile

from .manifest import load_manifest
from .normalizers import normalize_format_robustness_manifest

REQUIRED_FORMATS = ("pdf", "word", "pptx", "excel", "csv", "markdown", "text")


@dataclass(frozen=True, slots=True)
class FormatSmokeSpec:
    format_type: str
    suffix: str
    file_name: str
    question: str
    answer: str
    body: str


@dataclass(slots=True)
class IndexedDocument:
    document_id: str
    format_type: str
    document_path: str
    status: str
    text_chars: int = 0
    text: str = ""
    failure_type: str | None = None
    error: str | None = None

    def to_report(self) -> dict[str, Any]:
        payload = asdict(self)
        payload.pop("text", None)
        return payload


FORMAT_SMOKE_SPECS = (
    FormatSmokeSpec(
        format_type="pdf",
        suffix=".pdf",
        file_name="format_smoke_pdf.pdf",
        question="What is the PDF smoke answer?",
        answer="PDF smoke answer",
        body="This PDF fixture contains the PDF smoke answer for indexing.",
    ),
    FormatSmokeSpec(
        format_type="word",
        suffix=".docx",
        file_name="format_smoke_word.docx",
        question="What is the Word smoke answer?",
        answer="Word smoke answer",
        body="This Word fixture contains the Word smoke answer for indexing.",
    ),
    FormatSmokeSpec(
        format_type="pptx",
        suffix=".pptx",
        file_name="format_smoke_pptx.pptx",
        question="What is the PPTX smoke answer?",
        answer="PPTX smoke answer",
        body="This PPTX fixture contains the PPTX smoke answer for indexing.",
    ),
    FormatSmokeSpec(
        format_type="excel",
        suffix=".xlsx",
        file_name="format_smoke_excel.xlsx",
        question="What is the Excel smoke answer?",
        answer="Excel smoke answer",
        body="This Excel fixture contains the Excel smoke answer for indexing.",
    ),
    FormatSmokeSpec(
        format_type="csv",
        suffix=".csv",
        file_name="format_smoke_csv.csv",
        question="What is the CSV smoke answer?",
        answer="CSV smoke answer",
        body="This CSV fixture contains the CSV smoke answer for indexing.",
    ),
    FormatSmokeSpec(
        format_type="markdown",
        suffix=".md",
        file_name="format_smoke_markdown.md",
        question="What is the Markdown smoke answer?",
        answer="Markdown smoke answer",
        body="# Smoke\n\nThis Markdown fixture contains the Markdown smoke answer.",
    ),
    FormatSmokeSpec(
        format_type="text",
        suffix=".txt",
        file_name="format_smoke_text.txt",
        question="What is the text smoke answer?",
        answer="Text smoke answer",
        body="This text fixture contains the Text smoke answer for indexing.",
    ),
)

_EXTRACTORS: dict[str, Callable[[Path], str]] = {
    ".pdf": lambda path: _extract_pdf_text(path),
    ".docx": lambda path: _extract_docx_text(path),
    ".pptx": lambda path: _extract_pptx_text(path),
    ".xlsx": lambda path: _extract_xlsx_text(path),
    ".csv": lambda path: _extract_csv_text(path),
    ".md": lambda path: path.read_text(encoding="utf-8"),
    ".markdown": lambda path: path.read_text(encoding="utf-8"),
    ".txt": lambda path: path.read_text(encoding="utf-8"),
}


def build_format_smoke_fixtures(
    source_dir: str | Path,
    manifest_path: str | Path,
) -> Path:
    source_path = Path(source_dir)
    source_path.mkdir(parents=True, exist_ok=True)
    for spec in FORMAT_SMOKE_SPECS:
        format_dir = source_path / spec.format_type
        format_dir.mkdir(parents=True, exist_ok=True)
        document_path = format_dir / f"1_{spec.file_name}"
        _write_fixture_document(document_path, spec)
        _write_fixture_metadata(format_dir / "1_metadata.json", spec)
    return normalize_format_robustness_manifest(source_path, manifest_path)


def run_format_smoke_harness(manifest_path: str | Path) -> dict[str, Any]:
    bundle = load_manifest(manifest_path)
    indexed = {
        document_id: _index_document(document)
        for document_id, document in bundle.documents.items()
    }
    queries = [
        _query_example(example, indexed.get(example.document_id))
        for example in bundle.examples
    ]
    indexing_rows = [item.to_report() for item in indexed.values()]
    return {
        "schema_version": 1,
        "dataset_name": bundle.dataset_name,
        "manifest_path": str(Path(manifest_path)),
        "required_formats": list(REQUIRED_FORMATS),
        "overall_status": _overall_status(indexing_rows, queries),
        "num_documents": len(bundle.documents),
        "num_examples": len(bundle.examples),
        "indexing": indexing_rows,
        "queries": queries,
        "format_summary": _format_summary(indexing_rows, queries),
    }


def write_format_smoke_report(report: dict[str, Any], output_path: str | Path) -> Path:
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(report, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return path


def _write_fixture_metadata(path: Path, spec: FormatSmokeSpec) -> None:
    path.write_text(
        json.dumps(
            {
                "file_name": spec.file_name,
                "questions": [{"question": spec.question, "answer": spec.answer}],
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )


def _write_fixture_document(path: Path, spec: FormatSmokeSpec) -> None:
    if spec.suffix == ".pdf":
        _write_pdf(path, spec.body)
    elif spec.suffix == ".docx":
        _write_docx(path, spec.body)
    elif spec.suffix == ".pptx":
        _write_pptx(path, spec.body)
    elif spec.suffix == ".xlsx":
        _write_xlsx(path, spec.body)
    elif spec.suffix == ".csv":
        _write_csv(path, spec.body)
    else:
        path.write_text(spec.body + "\n", encoding="utf-8")


def _write_pdf(path: Path, text: str) -> None:
    import fitz

    document = fitz.open()
    try:
        page = document.new_page()
        page.insert_text((72, 72), text)
        document.save(path)
    finally:
        document.close()


def _write_docx(path: Path, text: str) -> None:
    from docx import Document

    document = Document()
    document.add_paragraph(text)
    document.save(path)


def _write_pptx(path: Path, text: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    box.text_frame.text = text
    presentation.save(path)


def _write_xlsx(path: Path, text: str) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Smoke"
    worksheet["A1"] = "content"
    worksheet["A2"] = text
    workbook.save(path)


def _write_csv(path: Path, text: str) -> None:
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerow(["content"])
        writer.writerow([text])


def _index_document(document: Any) -> IndexedDocument:
    path = Path(document.path)
    if not path.exists():
        return _failed_index(document, "missing_document", f"{path} does not exist")
    extractor = _EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        return _failed_index(
            document,
            "unsupported_format",
            f"Unsupported suffix: {path.suffix.lower()}",
        )
    try:
        text = extractor(path).strip()
    except (
        BadZipFile,
        ImportError,
        KeyError,
        OSError,
        RuntimeError,
        ValueError,
    ) as exc:
        return _failed_index(document, "index_error", str(exc))
    if not text:
        return _failed_index(document, "empty_index", "No text was extracted.")
    return IndexedDocument(
        document_id=document.document_id,
        format_type=document.format_type,
        document_path=str(path),
        status="pass",
        text_chars=len(text),
        text=text,
    )


def _failed_index(document: Any, failure_type: str, error: str) -> IndexedDocument:
    return IndexedDocument(
        document_id=document.document_id,
        format_type=document.format_type,
        document_path=str(document.path),
        status="fail",
        failure_type=failure_type,
        error=error,
    )


def _query_example(example: Any, indexed: IndexedDocument | None) -> dict[str, Any]:
    format_type = indexed.format_type if indexed is not None else "unknown"
    base = {
        "example_id": example.example_id,
        "document_id": example.document_id,
        "format_type": format_type,
        "question": example.question,
        "expected_answers": list(example.answers),
    }
    if indexed is None:
        return {**base, "status": "fail", "failure_type": "missing_index"}
    if indexed.status != "pass":
        return {**base, "status": "fail", "failure_type": indexed.failure_type}
    for answer in example.answers:
        if _contains_answer(indexed.text, answer):
            return {**base, "status": "pass", "predicted_answer": answer}
    return {**base, "status": "fail", "failure_type": "answer_not_indexed"}


def _contains_answer(text: str, answer: str) -> bool:
    return str(answer or "").casefold() in str(text or "").casefold()


def _overall_status(
    indexing_rows: list[dict[str, Any]],
    query_rows: list[dict[str, Any]],
) -> str:
    rows = [*indexing_rows, *query_rows]
    return "pass" if rows and all(item["status"] == "pass" for item in rows) else "fail"


def _format_summary(
    indexing_rows: list[dict[str, Any]],
    query_rows: list[dict[str, Any]],
) -> dict[str, dict[str, Any]]:
    summary: dict[str, dict[str, Any]] = {
        format_type: _new_summary(format_type) for format_type in REQUIRED_FORMATS
    }
    for row in indexing_rows:
        format_type = str(row["format_type"])
        item = summary.setdefault(format_type, _new_summary(format_type))
        item["documents"] += 1
        if row["status"] == "pass":
            item["indexed"] += 1
        else:
            _count_failure(item["failures"], row.get("failure_type"))
    for row in query_rows:
        format_type = str(row["format_type"])
        item = summary.setdefault(format_type, _new_summary(format_type))
        item["examples"] += 1
        if row["status"] == "pass":
            item["query_passed"] += 1
        else:
            _count_failure(item["failures"], row.get("failure_type"))
    for item in summary.values():
        if (
            item["documents"] > 0
            and item["examples"] > 0
            and item["documents"] == item["indexed"]
            and item["examples"] == item["query_passed"]
        ):
            item["status"] = "pass"
    return summary


def _new_summary(format_type: str) -> dict[str, Any]:
    return {
        "format_type": format_type,
        "documents": 0,
        "examples": 0,
        "indexed": 0,
        "query_passed": 0,
        "failures": {},
        "status": "fail",
    }


def _count_failure(failures: dict[str, int], failure_type: Any) -> None:
    key = str(failure_type or "unknown")
    failures[key] = failures.get(key, 0) + 1


def _extract_pdf_text(path: Path) -> str:
    import fitz

    document = fitz.open(path)
    try:
        return "\n".join(page.get_text() for page in document)
    finally:
        document.close()


def _extract_docx_text(path: Path) -> str:
    from docx import Document

    document = Document(path)
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def _extract_xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    try:
        parts: list[str] = []
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows(values_only=True):
                for value in row:
                    if value not in (None, ""):
                        parts.append(str(value))
        return "\n".join(parts)
    finally:
        workbook.close()


def _extract_csv_text(path: Path) -> str:
    with path.open(encoding="utf-8", newline="") as handle:
        return "\n".join(
            " ".join(cell for cell in row if cell) for row in csv.reader(handle)
        )
